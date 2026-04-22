#!/usr/bin/env python3
"""
2026 W6 (2/3-2/9) 財金股市週報 — PPT 生成腳本
風格：Corporate Blue & White Business Theme
內容：本週市場總覽、重大事件、總體經濟、產業趨勢、投資策略
資料來源：財金號角 SRT (2/5, 2/6, 2/9)、股癌投資 EP633
"""

import os
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from io import BytesIO
from math import pi

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ============================================================
# Config
# ============================================================

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(SCRIPT_DIR, '..', '..', '..'))
OUTPUT_DIR = os.path.join(PROJECT_ROOT, 'output', 'reports')
TEMP_DIR = os.path.join(PROJECT_ROOT, 'output', 'temp')
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(TEMP_DIR, exist_ok=True)
OUTPUT_PPTX = os.path.join(OUTPUT_DIR, 'weekly_report.pptx')

# Colors — Corporate Blue Theme
CORP_BLUE = RGBColor(0x2B, 0x57, 0x9A)
CORP_BLUE_DARK = RGBColor(0x1A, 0x3A, 0x6C)
CORP_BLUE_LIGHT = RGBColor(0x4A, 0x86, 0xC8)
ACCENT_BLUE = RGBColor(0x5B, 0xA0, 0xD6)
SKY_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
CHARCOAL = RGBColor(0x2D, 0x2D, 0x2D)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x4D)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
BULL_GREEN = RGBColor(0x26, 0xA6, 0x9A)
BEAR_RED = RGBColor(0xEF, 0x53, 0x50)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_CN = 'Microsoft JhengHei'
FONT_EN = 'Arial'
TOTAL_PAGES = 20

# ============================================================
# Matplotlib Setup
# ============================================================

def setup_matplotlib():
    plt.rcParams['font.family'] = ['Heiti TC', 'Arial Unicode MS', 'Hiragino Sans', 'sans-serif']
    plt.rcParams['axes.unicode_minus'] = False
    plt.rcParams['figure.dpi'] = 150
    plt.rcParams['savefig.dpi'] = 250
    plt.rcParams['font.size'] = 14
    plt.rcParams['axes.labelsize'] = 15
    plt.rcParams['axes.titlesize'] = 18


def fig_to_image_stream(fig, dpi=250):
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor='white', edgecolor='none', pad_inches=0.15)
    plt.close(fig)
    buf.seek(0)
    return buf


# ============================================================
# Helper functions
# ============================================================

def _add_bg_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        _set_shape_alpha(shape, alpha)
    return shape


def _set_shape_alpha(shape, alpha_pct):
    fill_elem = shape.fill._fill
    solid = fill_elem.find(qn('a:solidFill'))
    if solid is not None:
        srgb = solid.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_elem = srgb.find(qn('a:alpha'))
            if alpha_elem is None:
                from lxml import etree
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
            alpha_elem.set('val', str(int((100 - alpha_pct) * 1000)))


def _add_text_box(slide, left, top, width, height, text, font_size=14,
                  font_color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name=FONT_CN, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    sp = p._pPr
    if sp is None:
        from lxml import etree
        sp = etree.SubElement(p._p, qn('a:pPr'))
    lnSpc = sp.find(qn('a:lnSpc'))
    if lnSpc is None:
        from lxml import etree
        lnSpc = etree.SubElement(sp, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))
    return txBox


def _add_multiline_text_box(slide, left, top, width, height, lines,
                            font_size=14, font_color=DARK_GRAY,
                            font_name=FONT_CN, line_spacing=1.5,
                            alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_data, dict):
            p.text = line_data.get('text', '')
            p.font.size = Pt(line_data.get('size', font_size))
            p.font.color.rgb = line_data.get('color', font_color)
            p.font.bold = line_data.get('bold', False)
            p.font.name = line_data.get('font', font_name)
        else:
            p.text = str(line_data)
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(2)
    return txBox


def _add_circle(slide, left, top, size, color, text='', font_size=14,
                font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = FONT_EN
        p.alignment = PP_ALIGN.CENTER
    return shape


def _add_rounded_rect(slide, left, top, width, height, color, text='',
                      font_size=13, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = FONT_CN
        p.alignment = PP_ALIGN.CENTER
    return shape


def _add_header_bar(slide, title_text, subtitle_text=''):
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), CORP_BLUE)
    _add_circle(slide, Inches(0.3), Inches(0.15), Inches(0.8), CORP_BLUE_DARK)
    _add_text_box(slide, Inches(1.3), Inches(0.15), Inches(10), Inches(0.55),
                  title_text, font_size=28, font_color=WHITE, bold=True)
    if subtitle_text:
        _add_text_box(slide, Inches(1.3), Inches(0.65), Inches(10), Inches(0.4),
                      subtitle_text, font_size=15, font_color=SKY_BLUE)
    _add_text_box(slide, Inches(10.5), Inches(0.15), Inches(2.5), Inches(0.4),
                  '財金週報', font_size=16, font_color=WHITE, bold=True,
                  alignment=PP_ALIGN.RIGHT)
    _add_text_box(slide, Inches(10.5), Inches(0.55), Inches(2.5), Inches(0.35),
                  '2026.02.03 - 02.09', font_size=12, font_color=SKY_BLUE,
                  alignment=PP_ALIGN.RIGHT)


def _add_footer(slide, page_num, total_pages=TOTAL_PAGES):
    _add_bg_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), LIGHT_GRAY)
    _add_text_box(slide, Inches(0.5), Inches(7.15), Inches(5), Inches(0.3),
                  '2026 W6 財金股市週報 — 全面市場分析', font_size=10,
                  font_color=MID_GRAY)
    _add_text_box(slide, Inches(10), Inches(7.15), Inches(3), Inches(0.3),
                  f'{page_num} / {total_pages}', font_size=10,
                  font_color=MID_GRAY, alignment=PP_ALIGN.RIGHT, font_name=FONT_EN)


def _add_section_divider(slide, num, title, subtitle, detail=''):
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CORP_BLUE)
    _add_circle(slide, Inches(1.5), Inches(2), Inches(2.5), WHITE, num,
                font_size=48, font_color=CORP_BLUE)
    _add_text_box(slide, Inches(5), Inches(2.2), Inches(7), Inches(1),
                  title, font_size=40, font_color=WHITE, bold=True, font_name=FONT_EN)
    _add_text_box(slide, Inches(5), Inches(3.3), Inches(7), Inches(0.6),
                  subtitle, font_size=18, font_color=SKY_BLUE)
    _add_bg_rect(slide, Inches(5), Inches(4.1), Inches(4), Inches(0.04), WHITE)
    if detail:
        _add_text_box(slide, Inches(5), Inches(4.4), Inches(7), Inches(0.5),
                      detail, font_size=14, font_color=SKY_BLUE)


# ============================================================
# Chart Generation
# ============================================================

def gen_charts():
    charts = {}
    blue_hex = '#2B579A'
    blue_light_hex = '#4A86C8'
    accent_hex = '#5BA0D6'
    green_hex = '#27AE60'
    red_hex = '#E84D4D'
    orange_hex = '#F59E0B'
    purple_hex = '#8E44AD'

    # ── 1. Weekly US Indices Performance ──
    fig, ax = plt.subplots(figsize=(11, 3.7))
    indices = ['道瓊工業\n50,115', '標普 500\n6,026', '納斯達克\n19,523', '費城半導體\n4,856']
    weekly_chg = [4.67, -0.24, -1.63, -3.12]
    colors = [green_hex if v > 0 else red_hex for v in weekly_chg]
    bars = ax.barh(indices, weekly_chg, color=colors, height=0.5, edgecolor='white', linewidth=1.5)
    ax.axvline(x=0, color='#ccc', linewidth=1.2, linestyle='-')
    for bar, val in zip(bars, weekly_chg):
        x_pos = bar.get_width() + (0.15 if val >= 0 else -0.15)
        ha = 'left' if val >= 0 else 'right'
        ax.text(x_pos, bar.get_y() + bar.get_height() / 2,
                f'{val:+.2f}%', va='center', ha=ha, fontsize=16,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xlim(-5, 6)
    ax.set_title('美股四大指數本週表現（2026/2/3-2/7）', fontsize=18,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.tick_params(axis='both', labelsize=14)
    fig.tight_layout()
    charts['weekly_us'] = fig_to_image_stream(fig)

    # ── 2. Taiwan & Asia Weekly ──
    fig, ax = plt.subplots(figsize=(11, 3.5))
    asia_names = ['台灣加權\n23,611', '日經 225\n39,500', '恆生指數\n20,890',
                  '上證指數\n3,280', '韓國 KOSPI\n2,520']
    asia_chg = [2.48, 3.15, 1.82, 0.95, 1.35]
    asia_colors = [green_hex if v > 0 else red_hex for v in asia_chg]
    bars = ax.bar(range(len(asia_names)), asia_chg, color=asia_colors,
                  width=0.5, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, asia_chg)):
        y_pos = val + 0.12
        ax.text(i, y_pos, f'+{val:.2f}%', ha='center', fontsize=14,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(asia_names)))
    ax.set_xticklabels(asia_names, fontsize=12)
    ax.set_title('亞股本週表現（2026/2/3-2/9）', fontsize=18,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.axhline(y=0, color='#ccc', linewidth=1)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(0, 4.5)
    fig.tight_layout()
    charts['asia_weekly'] = fig_to_image_stream(fig)

    # ── 3. Bitcoin Crash ──
    fig, ax = plt.subplots(figsize=(7.5, 3.5))
    btc_months = ['2025/10', '11', '12', '2026/1', '2/3', '2/5', '2/6', '2/7']
    btc_prices = [95000, 88000, 82000, 78000, 72000, 68000, 63000, 65000]
    ax.plot(btc_months, btc_prices, 'o-', color=red_hex, linewidth=3,
            markersize=10, zorder=5)
    ax.fill_between(btc_months, btc_prices, alpha=0.15, color=red_hex)
    ax.axhline(y=90000, color=orange_hex, linewidth=1.5, linestyle='--',
               label='川普就任價位', alpha=0.8)
    ax.axhline(y=50000, color='#999', linewidth=1, linestyle=':', label='200 週均線')
    ax.legend(fontsize=13, loc='upper right', framealpha=0.9)
    ax.set_ylabel('比特幣價格 (USD)', fontsize=14, color=blue_hex)
    ax.set_title('比特幣崩跌：跌穿 $63,000，川普行情全數蒸發', fontsize=17,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(40000, 100000)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f'${x/1000:.0f}K'))
    fig.tight_layout()
    charts['btc_crash'] = fig_to_image_stream(fig)

    # ── 4. ISM PMI Recovery ──
    fig, ax = plt.subplots(figsize=(8, 3.5))
    pmi_years = ['2022', '2023\nQ1', 'Q2', 'Q3', 'Q4',
                 '2024\nQ1', 'Q2', 'Q3', 'Q4',
                 '2025\nQ1', 'Q2', 'Q3', 'Q4', '2026\n1月']
    pmi_values = [52.8, 48.5, 47.2, 46.5, 47.8,
                  48.2, 47.5, 48.8, 49.3,
                  49.8, 48.5, 49.2, 50.5, 52.6]
    colors_pmi = [green_hex if v >= 50 else red_hex for v in pmi_values]
    bars = ax.bar(range(len(pmi_years)), pmi_values, color=colors_pmi,
                  width=0.6, edgecolor='white', linewidth=1)
    ax.axhline(y=50, color=blue_hex, linewidth=2, linestyle='--',
               label='榮枯線 (50)', alpha=0.8)
    ax.text(len(pmi_years) - 1, 52.6 + 0.8, '52.6', ha='center',
            fontsize=15, fontweight='bold', color=green_hex)
    ax.set_xticks(range(len(pmi_years)))
    ax.set_xticklabels(pmi_years, fontsize=10.5)
    ax.legend(fontsize=14, loc='lower right')
    ax.set_ylabel('ISM PMI', fontsize=14, color=blue_hex)
    ax.set_title('美國 ISM 製造業 PMI：突破 3 年庫存調整', fontsize=18,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.set_ylim(43, 56)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    fig.tight_layout()
    charts['ism_pmi'] = fig_to_image_stream(fig)

    # ── 5. Value vs Growth ──
    fig, ax = plt.subplots(figsize=(7.5, 3.5))
    months = ['2025/10', '11', '12', '2026/1', '2/7']
    value_returns = [0.5, 2.1, 3.8, 5.5, 7.2]
    growth_returns = [1.2, -0.3, -1.8, -2.5, -3.0]
    ax.plot(months, value_returns, 'o-', color=blue_hex, linewidth=3,
            markersize=10, label='羅素1000 價值股', zorder=5)
    ax.plot(months, growth_returns, 's-', color=red_hex, linewidth=3,
            markersize=10, label='羅素1000 成長股', zorder=5)
    ax.fill_between(months, value_returns, alpha=0.15, color=blue_hex)
    ax.fill_between(months, growth_returns, alpha=0.15, color=red_hex)
    ax.axhline(y=0, color='#999', linewidth=1, linestyle='--')
    ax.legend(fontsize=15, loc='upper left', framealpha=0.9)
    ax.set_ylabel('累積報酬率 (%)', fontsize=14, color=blue_hex)
    ax.set_title('價值股 vs 成長股：2022 年以來最大分歧', fontsize=18,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    fig.tight_layout()
    charts['value_growth'] = fig_to_image_stream(fig)

    # ── 6. Tech Capex War ──
    fig, ax = plt.subplots(figsize=(10, 3.8))
    companies = ['Google\n$1,800B', 'Amazon\n$2,000B', 'Meta\n$650B',
                 'Microsoft\n$800B', 'Apple\n$200B']
    capex_2025 = [800, 1300, 450, 550, 120]
    capex_2026 = [1800, 2000, 650, 800, 200]
    x = np.arange(len(companies))
    w = 0.3
    bars1 = ax.bar(x - w/2, capex_2025, w, color=accent_hex, label='2025 Capex',
                   edgecolor='white', linewidth=1.5)
    bars2 = ax.bar(x + w/2, capex_2026, w, color=blue_hex, label='2026E Capex',
                   edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars2, capex_2026):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 30,
                f'${val}B', ha='center', fontsize=12, fontweight='bold', color=blue_hex)
    ax.set_xticks(x)
    ax.set_xticklabels(companies, fontsize=12)
    ax.set_ylabel('資本支出 (億美元)', fontsize=14, color=blue_hex)
    ax.set_title('科技巨頭 2026 年資本支出軍備競賽', fontsize=18,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.legend(fontsize=13, loc='upper right')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(0, 2400)
    fig.tight_layout()
    charts['capex_war'] = fig_to_image_stream(fig)

    # ── 7. Sector Radar ──
    fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
    categories = ['低軌衛星', 'AI 散熱', '國防軍工', '傳產價值',
                  'AI 軟體', 'AI 硬體', '能源電力']
    N = len(categories)
    current = [9, 8, 7, 8, 4, 6, 7]
    future = [9, 9, 8, 7, 6, 8, 8]
    angles = [n / float(N) * 2 * pi for n in range(N)]
    angles += angles[:1]
    current += current[:1]
    future += future[:1]
    ax.plot(angles, current, 'o-', color=blue_hex, linewidth=2.5,
            label='本週動能', markersize=8)
    ax.fill(angles, current, alpha=0.15, color=blue_hex)
    ax.plot(angles, future, 's--', color=orange_hex, linewidth=2.5,
            label='未來 1-3 月展望', markersize=8)
    ax.fill(angles, future, alpha=0.1, color=orange_hex)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories, fontsize=14, fontweight='bold')
    ax.set_ylim(0, 10)
    ax.set_yticks([2, 4, 6, 8, 10])
    ax.set_yticklabels(['2', '4', '6', '8', '10'], fontsize=11, color='gray')
    ax.legend(fontsize=14, loc='upper right', bbox_to_anchor=(1.3, 1.1))
    ax.set_title('產業動能雷達圖', fontsize=18,
                 fontweight='bold', color=blue_hex, pad=30)
    fig.tight_layout()
    charts['sector_radar'] = fig_to_image_stream(fig)

    # ── 8. Asset Drawdown ──
    fig, ax = plt.subplots(figsize=(10, 3.5))
    assets = ['比特幣', '白銀\n(盤中)', '軟體股\nETF', '費半\nYTD', '黃金']
    drawdowns = [-30, -40, -20, -3.12, -12]
    dd_colors = ['#C0392B', '#E74C3C', red_hex, orange_hex, '#F39C12']
    bars = ax.barh(assets, drawdowns, color=dd_colors, height=0.5,
                   edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars, drawdowns):
        ax.text(bar.get_width() - 1, bar.get_y() + bar.get_height() / 2,
                f'{val}%', va='center', ha='right', fontsize=15,
                fontweight='bold', color='white')
    ax.axvline(x=0, color='#ccc', linewidth=1.2)
    ax.set_title('本週各類資產回檔幅度', fontsize=18,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.set_xlim(-50, 5)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.tick_params(axis='both', labelsize=13)
    fig.tight_layout()
    charts['asset_drawdown'] = fig_to_image_stream(fig)

    # ── 9. Allocation Pie ──
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(8.5, 3.5))
    labels_c = ['科技龍頭 ETF', '價值/傳產', '債券/現金', '黃金/商品']
    sizes_c = [25, 35, 30, 10]
    colors_c = [blue_hex, green_hex, '#95a5a6', orange_hex]
    wedges, texts, autotexts = ax1.pie(
        sizes_c, labels=labels_c, autopct='%1.0f%%', colors=colors_c,
        startangle=90, pctdistance=0.75, textprops={'fontsize': 13})
    for at in autotexts:
        at.set_fontweight('bold')
        at.set_color('white')
    ax1.set_title('穩健型配置', fontsize=16, fontweight='bold', color=blue_hex, pad=15)

    labels_g = ['AI 硬體/半導體', '低軌衛星/散熱', '傳產循環', '軟體(低接)', '現金']
    sizes_g = [30, 20, 25, 15, 10]
    colors_g = [blue_hex, purple_hex, green_hex, accent_hex, '#95a5a6']
    wedges2, texts2, autotexts2 = ax2.pie(
        sizes_g, labels=labels_g, autopct='%1.0f%%', colors=colors_g,
        startangle=90, pctdistance=0.75, textprops={'fontsize': 12})
    for at in autotexts2:
        at.set_fontweight('bold')
        at.set_color('white')
    ax2.set_title('積極型配置', fontsize=16, fontweight='bold', color=blue_hex, pad=15)
    fig.tight_layout()
    charts['allocation'] = fig_to_image_stream(fig)

    # ── 10. Japan Defense Spending ──
    fig, ax = plt.subplots(figsize=(7, 3.5))
    pm_names = ['安倍', '菅義偉', '岸田文雄', '石破茂', '高市早苗\n(2026E)']
    defense_pct = [5.3, 5.0, 7.0, 7.5, 8.5]
    bar_colors_jp = [accent_hex, accent_hex, blue_light_hex, blue_light_hex, blue_hex]
    bars = ax.bar(pm_names, defense_pct, color=bar_colors_jp, width=0.5,
                  edgecolor='white', linewidth=2)
    for bar, val in zip(bars, defense_pct):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.15,
                f'{val}%', ha='center', fontsize=14, fontweight='bold', color=blue_hex)
    ax.set_ylabel('國防支出占一般預算比例 (%)', fontsize=13, color=blue_hex)
    ax.set_title('日本歷任首相國防支出預算占比', fontsize=17,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(0, 10)
    fig.tight_layout()
    charts['japan_defense'] = fig_to_image_stream(fig)

    return charts


# ============================================================
# Slide Builders
# ============================================================

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CORP_BLUE)
    _add_circle(slide, Inches(10.5), Inches(-0.5), Inches(3), CORP_BLUE_DARK)
    _add_circle(slide, Inches(11.5), Inches(5), Inches(2.5), CORP_BLUE_DARK)
    _add_circle(slide, Inches(-0.8), Inches(5.5), Inches(2), CORP_BLUE_DARK)
    _add_rounded_rect(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(5), WHITE)
    _add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
                  '2026 W6 財金股市週報', font_size=44, font_color=CORP_BLUE,
                  bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.8),
                  '全面市場分析 ‧ 總經觀察 ‧ 產業趨勢 ‧ 投資策略',
                  font_size=24, font_color=CORP_BLUE_DARK,
                  bold=True, alignment=PP_ALIGN.CENTER)
    _add_bg_rect(slide, Inches(4.5), Inches(3.5), Inches(4.3), Inches(0.05), CORP_BLUE)
    _add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.5),
                  '川普推文道瓊創歷史新高 50,115 | 比特幣崩跌 30% | 日本大選高市完全執政',
                  font_size=14, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
                  'PMI 52.6 重返擴張 | 低軌衛星爆發 | 價值股創歷史新高',
                  font_size=14, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.4),
                  '2026 年 2 月 3 日 — 2 月 9 日', font_size=12,
                  font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(4), Inches(6.5), Inches(5.3), Inches(0.5),
                  '財金號角 ‧ 股癌投資 — 綜合分析', font_size=18,
                  font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_02_contents(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '目錄 CONTENTS', '本週全方位市場觀點與投資建議')
    _add_footer(slide, 2)

    sections = [
        {'num': '01', 'title': '市場總覽', 'desc': '美股週表現\n台股與亞股動態', 'color': CORP_BLUE},
        {'num': '02', 'title': '重大事件', 'desc': '川普推文拉盤\n日本大選\n比特幣崩跌', 'color': CORP_BLUE_LIGHT},
        {'num': '03', 'title': '總體經濟', 'desc': 'PMI 突破擴張\n價值 vs 成長\nAI 就業衝擊', 'color': ACCENT_BLUE},
        {'num': '04', 'title': '產業趨勢', 'desc': '低軌衛星爆發\n散熱族群突破\n資本支出軍備賽', 'color': ACCENT_PURPLE},
    ]

    start_x = Inches(0.8)
    for i, sec in enumerate(sections):
        x = start_x + i * Inches(3.2)
        y = Inches(2.0)
        _add_rounded_rect(slide, x, y, Inches(2.9), Inches(4.3), WHITE)
        _add_circle(slide, x + Inches(0.85), y + Inches(0.3), Inches(1.1),
                    sec['color'], sec['num'], font_size=28)
        _add_text_box(slide, x + Inches(0.1), y + Inches(1.6), Inches(2.7),
                      Inches(0.5), sec['title'], font_size=20,
                      font_color=CHARCOAL, bold=True, alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, x + Inches(0.2), y + Inches(2.2), Inches(2.5),
                      Inches(1.5), sec['desc'], font_size=13,
                      font_color=MID_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.5)
        _add_bg_rect(slide, x + Inches(0.8), y + Inches(4.0), Inches(1.3),
                     Inches(0.05), sec['color'])


def slide_03_part_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '01', 'PART ONE.', '本週市場總覽 — 美股、台股、亞股全景',
                         '道瓊 +4.67% 創歷史新高 | 費半 -3.12% | 台股 +2.48%')


def slide_04_us_weekly(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '美股四大指數本週表現', '2026/2/3-2/7 收盤數據')
    _add_footer(slide, 4)

    idx_data = [
        ('道瓊工業', '50,115', '+1,207 pt', '+4.67%', ACCENT_GREEN),
        ('標普 500', '6,026', '-15 pt', '-0.24%', ACCENT_RED),
        ('納斯達克', '19,523', '-323 pt', '-1.63%', ACCENT_RED),
        ('費城半導體', '4,856', '-156 pt', '-3.12%', ACCENT_RED),
    ]
    for i, (name, price, change, pct, color) in enumerate(idx_data):
        x = Inches(0.4) + i * Inches(3.2)
        _add_rounded_rect(slide, x, Inches(1.4), Inches(3.0), Inches(1.6), WHITE)
        _add_bg_rect(slide, x + Inches(0.05), Inches(1.4), Inches(2.9), Inches(0.06), color)
        _add_text_box(slide, x + Inches(0.2), Inches(1.55), Inches(2.6), Inches(0.35),
                      name, font_size=15, font_color=CHARCOAL, bold=True, alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(2.6), Inches(0.4),
                      price, font_size=24, font_color=CHARCOAL, bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
        _add_text_box(slide, x + Inches(0.2), Inches(2.35), Inches(2.6), Inches(0.35),
                      f'{change}（{pct}）', font_size=14, font_color=color,
                      bold=True, alignment=PP_ALIGN.CENTER)

    img = charts['weekly_us']
    slide.shapes.add_picture(img, Inches(1.2), Inches(3.2), Inches(11), Inches(3.7))


def slide_05_asia(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '台股與亞股本週動態', '日本大選行情爆發 | 台股強勢抗跌')
    _add_footer(slide, 5)

    img = charts['asia_weekly']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.2), Inches(1.3), Inches(4.8), Inches(3.5), WHITE)
    asia_lines = [
        {'text': '本週亞股重點', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  台股 32,611 點（+788, 2/9）', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  美股暴殺但台股跌幅極小，異常強勁', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  日經大漲 +1,000 點（2/9 開盤）', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  高市早苗大選壓倒性勝利推動', 'size': 12, 'color': DARK_GRAY},
        {'text': '  自民黨 352 席（過半僅需 233）', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  美中電話會談正面', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  川普 4 月訪中計劃確認', 'size': 12, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(8.4), Inches(1.45), Inches(4.4), Inches(3.2),
                            asia_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    bottom = [
        {'text': '台股觀察重點', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  台股 500 點以內都是正常波動，收 31,921 點（2/5）→ 32,611 點（2/9），整週反彈強勁', 'size': 14, 'color': DARK_GRAY},
        {'text': '  外資系統單尚未大規模啟動，台灣 PMI 57 領先全球，製造業上行循環支撐強勁', 'size': 14, 'color': DARK_GRAY},
        {'text': '  低軌衛星族群成避風港，大盤跌時跌少，漲時逆勢大漲 — 但盤好時需提防資金輪出', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), bottom)


def slide_06_part_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '02', 'PART TWO.', '本週重大事件 — 牽動全球市場走向',
                         '川普推文 | 日本大選 | 比特幣崩跌 | 資產去槓桿')


def slide_07_trump_china(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '川普推文拉盤 & 美中關係破冰', '道瓊 +1,207 點創歷史新高 50,115')
    _add_footer(slide, 7)

    # Trump event card
    _add_rounded_rect(slide, Inches(0.3), Inches(1.4), Inches(6.2), Inches(3.2), WHITE)
    trump_lines = [
        {'text': '川普推文效應', 'size': 18, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  週五（2/7）川普發文慶祝道瓊突破 5 萬點', 'size': 14, 'color': CHARCOAL},
        {'text': '  道瓊指數大漲 +1,207 點，收 50,115 點（歷史新高）', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  本輪科技股中期回檔似乎告一段落', 'size': 14, 'color': CHARCOAL},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  市場解讀：政治面消息產生大幅轉捩', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  但須觀察後續 — 1 月非農延至週三公佈', 'size': 13, 'color': MID_GRAY},
        {'text': '  值缺數持續下滑的效果值得關注', 'size': 13, 'color': MID_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(1.55), Inches(5.6), Inches(3.0),
                            trump_lines)

    # US-China card
    _add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(6.2), Inches(3.2), WHITE)
    china_lines = [
        {'text': '美中關係破冰', 'size': 18, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  川普與習近平進行「精彩且深入」的電話會談', 'size': 14, 'color': CHARCOAL},
        {'text': '  確認 4 月份川普將訪中', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  討論議題：貿易、軍事、台灣、俄烏、伊朗', 'size': 13, 'color': CHARCOAL},
        {'text': '  習近平提升大豆採購至 2,000 萬噸', 'size': 13, 'color': CHARCOAL},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  全球外交窗口打開', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  李在明訪陸、加拿大訪陸、英國開啟中英黃金時代', 'size': 12, 'color': MID_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(7.1), Inches(1.55), Inches(5.6), Inches(3.0),
                            china_lines)

    # Bottom
    _add_rounded_rect(slide, Inches(0.3), Inches(4.9), Inches(12.7), Inches(1.9), CORP_BLUE)
    bottom = [
        {'text': '市場影響評估', 'size': 18, 'color': WHITE, 'bold': True},
        {'text': '', 'size': 5, 'color': WHITE},
        {'text': '  短期利多：川普政治操盤穩定股市信心，道瓊創歷史新高帶動傳產與價值股資金流入', 'size': 15, 'color': WHITE},
        {'text': '  中期觀察：美中關係改善可降低全球貿易不確定性，有利製造業上行循環延續', 'size': 15, 'color': SKY_BLUE},
        {'text': '  風險提醒：每次股災後政治不確定性縮小 → 股市反彈 → 需觀察反彈力道是否持續', 'size': 15, 'color': SKY_BLUE},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.05), Inches(12), Inches(1.5), bottom)


def slide_08_japan(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '日本大選：高市早苗完全執政', '自民黨 352 席壓倒性勝利 — 國防、半導體、核能政策加速')
    _add_footer(slide, 8)

    # Defense chart
    img = charts['japan_defense']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7), Inches(3.5))

    # Policy impact
    _add_rounded_rect(slide, Inches(7.6), Inches(1.3), Inches(5.4), Inches(3.5), WHITE)
    jp_lines = [
        {'text': '高市政策三箭', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  國防擴張 → GDP 2% 目標', 'size': 15, 'color': ACCENT_RED, 'bold': True},
        {'text': '  三菱重工、川崎重工、三菱電機受惠', 'size': 12, 'color': DARK_GRAY},
        {'text': '  國防支出上看 10 兆日元', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  半導體戰略 → 台積電九州擴廠', 'size': 15, 'color': CORP_BLUE, 'bold': True},
        {'text': '  追加 $170 億投資，3 奈米製程同步落地', 'size': 12, 'color': DARK_GRAY},
        {'text': '  九州成為「半導體大洲」', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  食品消費稅 → 暫降零（2 年）', 'size': 15, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  財政赤字疑慮升高，日債壓力持續', 'size': 12, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(7.8), Inches(1.45), Inches(5.0), Inches(3.2),
                            jp_lines)

    # Bottom: Market Impact
    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    impact = [
        {'text': '市場影響：日股結構性分化', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  利多：軍工股、設備股、先進製造、電力、出口族群 → 政策題材+資金行情驅動日股創高', 'size': 14, 'color': ACCENT_GREEN},
        {'text': '  壓力：高股息防禦股失去殖利率優勢 | 10Y 國債 2.5%、超長期 4%+ → 股票相對吸引力受侵蝕', 'size': 14, 'color': ACCENT_RED},
        {'text': '  日元中期偏弱 | 財政赤字 GDP 250% | 日本核心通膨 2.4%（亞洲最高）→ 日銀政策兩難', 'size': 14, 'color': ACCENT_ORANGE},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), impact)


def slide_09_bitcoin(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '比特幣崩跌與風險資產大退潮', 'BTC 跌破 $63,000 | 白銀盤中 -17% | 軟體股 -20% YTD')
    _add_footer(slide, 9)

    # BTC chart
    img = charts['btc_crash']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.5))

    # Asset drawdown chart
    img2 = charts['asset_drawdown']
    slide.shapes.add_picture(img2, Inches(0.3), Inches(5.0), Inches(6.5), Inches(1.8))

    # Key points
    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(4.9), Inches(5.5), WHITE)
    btc_lines = [
        {'text': '風險資產全面退潮', 'size': 17, 'color': ACCENT_RED, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  比特幣 YTD -30%', 'size': 15, 'color': ACCENT_RED, 'bold': True},
        {'text': '  跌穿 $63K，川普任期漲幅全數蒸發', 'size': 12, 'color': DARK_GRAY},
        {'text': '  跌穿 100 週均線，200 週均線在 $60K', 'size': 12, 'color': DARK_GRAY},
        {'text': '  Gemini 交易所裁員 25% 關閉歐洲', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  白銀盤中暴跌 17%', 'size': 15, 'color': ACCENT_RED, 'bold': True},
        {'text': '  中國投機資金撤出', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  軟體股 ETF YTD -20%', 'size': 15, 'color': ACCENT_RED, 'bold': True},
        {'text': '  一個月內進入熊市', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  核心判斷', 'size': 15, 'color': CORP_BLUE, 'bold': True},
        {'text': '  不單純換股操作，是全面去槓桿', 'size': 13, 'color': CHARCOAL},
        {'text': '  動能交易策略踩踏 + 系統單獲利了結', 'size': 13, 'color': CHARCOAL},
        {'text': '  情緒未失控（VIX/VXN 未極端）', 'size': 13, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  乖離已嚴重超跌 → 左側佈局機會', 'size': 13, 'color': ACCENT_GREEN, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.5), Inches(5.2),
                            btc_lines)


def slide_10_part_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '03', 'PART THREE.', '總體經濟分析 — PMI、就業、資金輪動',
                         'ISM PMI 52.6 | 台灣 PMI 57 | 價值股領先成長股')


def slide_11_pmi(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '製造業 PMI 重返擴張：資金輪動啟動', 'ISM PMI 52.6 — 突破 3 年庫存調整')
    _add_footer(slide, 11)

    img = charts['ism_pmi']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(8), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.6), Inches(1.3), Inches(4.4), Inches(3.5), WHITE)
    pmi_lines = [
        {'text': '關鍵數據', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  美國 ISM PMI 52.6', 'size': 15, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  2022 年以來最大增幅', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  台灣 PMI 57（領先指標）', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  「賣鏟子的」帶貨潮率先啟動', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  新訂單 + 生產指數跳升', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  用戶端庫存下滑 → 補庫存需求', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  歷史規律', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  PMI < 50 = 股市買點', 'size': 13, 'color': CHARCOAL},
        {'text': '  突破 50 = 上行週期確認', 'size': 13, 'color': ACCENT_GREEN, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(8.8), Inches(1.45), Inches(4.0), Inches(3.2),
                            pmi_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    bottom = [
        {'text': '總經結論：庫存上行週期 + 中期回檔 = 最佳佈建時機', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  製造業上行循環已正式展開，科技股資金流往傳產實體製造業是結構性趨勢', 'size': 14, 'color': DARK_GRAY},
        {'text': '  歷史上庫存上行週期中的中期回檔，反而是最具吸引力的進場點', 'size': 14, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), bottom)


def slide_12_value_growth(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '價值股 vs 成長股：風格大切換', '2022 年以來最大分歧 — 資金從高估值轉往低估值')
    _add_footer(slide, 12)

    img = charts['value_growth']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.2), Inches(1.3), Inches(4.8), Inches(3.5), WHITE)
    val_lines = [
        {'text': '價值股代表（持續創高）', 'size': 16, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  Coca-Cola (KO)', 'size': 15, 'color': CHARCOAL, 'bold': True, 'font': FONT_EN},
        {'text': '  $60+ → $77，歷史新高', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  Walmart (WMT)', 'size': 15, 'color': CHARCOAL, 'bold': True, 'font': FONT_EN},
        {'text': '  2026 完全噴出，持續創高', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  道瓊指數', 'size': 15, 'color': CHARCOAL, 'bold': True},
        {'text': '  50,115 歷史新高，傳產推動', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 8, 'color': DARK_GRAY},
        {'text': '  風格切換，非系統性崩盤', 'size': 14, 'color': CORP_BLUE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(8.4), Inches(1.45), Inches(4.4), Inches(3.2),
                            val_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    insight = [
        {'text': '投資啟示', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  方向先確定，方法才有意義 — 確認自己是價值型還是成長型，擇一堅持', 'size': 14, 'color': DARK_GRAY},
        {'text': '  「我知道我是誰，但我不知道我要去哪裡」— 不要漲時追成長、跌時改價值、震盪做波段', 'size': 14, 'color': DARK_GRAY},
        {'text': '  景氣低基期亂買股票 > 景氣高基期認真研究財報 — 週期位置比個股選擇更重要', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), insight)


def slide_13_ai_jobs(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, 'AI 就業衝擊與軟體股重新定價', '1 月裁員 10.8 萬人（2009 年以來最嚴重）| 軟體股 Token 經濟新模式')
    _add_footer(slide, 13)

    # AI impact card
    _add_rounded_rect(slide, Inches(0.3), Inches(1.4), Inches(6.2), Inches(2.8), WHITE)
    ai_lines = [
        {'text': 'AI 對就業的雙面性', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  總體：AI 帶來生產力跳升，GDP 加速', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  個體：重複性高、可標準化的職位被替代', 'size': 14, 'color': ACCENT_RED, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  1 月裁員 10.8 萬人（2009 年以來最嚴重）', 'size': 13, 'color': CHARCOAL},
        {'text': '  科技業：調高 Capex + 同時裁員 + 投注 AI', 'size': 13, 'color': CHARCOAL},
        {'text': '  初領失業金 23.1 萬（高於預期）', 'size': 13, 'color': CHARCOAL},
        {'text': '  製造業職位空缺比 2023 年少 20 萬', 'size': 13, 'color': CHARCOAL},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(1.55), Inches(5.6), Inches(2.5),
                            ai_lines)

    # Software stock future
    _add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(6.2), Inches(2.8), WHITE)
    sw_lines = [
        {'text': '軟體股 Token 經濟新模式', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  A 劇本（悲觀）', 'size': 14, 'color': ACCENT_RED, 'bold': True},
        {'text': '  Token 成本自行吸收 → 毛利受擠壓', 'size': 12, 'color': DARK_GRAY},
        {'text': '  競爭對手用 Vibe Coding 產出類似產品', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  B 劇本（樂觀）', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  SEED + Token 雙軌收費 → Upsell', 'size': 12, 'color': DARK_GRAY},
        {'text': '  CSP Token 成本持續下降 → 毛利反升', 'size': 12, 'color': DARK_GRAY},
        {'text': '  B2B 軟體護城河高，不易被巨頭取代', 'size': 12, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(7.1), Inches(1.55), Inches(5.6), Inches(2.5),
                            sw_lines)

    # Bottom conclusion
    _add_rounded_rect(slide, Inches(0.3), Inches(4.5), Inches(12.7), Inches(2.3), WHITE)
    conclusion = [
        {'text': '本週判斷', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  Anthropic 推出 AI 法律軟體 → 市場解讀為「AI 正式跨進專業白領領域」→ 軟體股恐慌拋售', 'size': 14, 'color': DARK_GRAY},
        {'text': '  Cloudbot 等 Agent 工具崛起 → 但從 Prototype 到企業端 Workflow 整合仍有長路', 'size': 14, 'color': DARK_GRAY},
        {'text': '  2C 軟體被巨頭吃掉的風險高，但 2B 軟體護城河仍在 → 不應非黑即白看待', 'size': 14, 'color': DARK_GRAY},
        {'text': '  資安類股可能成為下一個受惠方向 — AI Agent 的權限與資料安全議題將大爆發', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  軟體股 RSI 已至 07-08 金融海嘯水位，乖離超跌但情緒未失控 → 具左側佈局價值', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(4.65), Inches(12), Inches(2.0),
                            conclusion)


def slide_14_part_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '04', 'PART FOUR.', '產業趨勢 — 值得投入的方向',
                         '低軌衛星 | AI 散熱 | 科技巨頭 Capex | 能源電力')


def slide_15_sectors(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '本週值得關注產業：低軌衛星 & AI 散熱', '2026-2027 絕地大爆發的兩大方向')
    _add_footer(slide, 15)

    # LEO Satellites
    _add_rounded_rect(slide, Inches(0.3), Inches(1.4), Inches(6.2), Inches(5.3), WHITE)
    leo_lines = [
        {'text': '低軌衛星（LEO Satellite）', 'size': 18, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  動能評級：★★★★★', 'size': 15, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '', 'size': 3, 'color': DARK_GRAY},
        {'text': '  市場地位', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  • 無庸置疑的強勢族群 — 大盤跌時避風港', 'size': 13, 'color': DARK_GRAY},
        {'text': '  • 大盤跌少，逆勢漲很兇', 'size': 13, 'color': DARK_GRAY},
        {'text': '  • 訂單數量驚人，供應鏈已全面確認', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 3, 'color': DARK_GRAY},
        {'text': '  供應鏈延伸', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  • PCB 基板：華通、耀華、星星、Meiko', 'size': 13, 'color': DARK_GRAY},
        {'text': '  • CCL 材料：M6 等級（非 AI 高階 M8/M9）', 'size': 13, 'color': DARK_GRAY},
        {'text': '  • 衣布（PCB 材料）：AI Server 排擠產能', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 3, 'color': DARK_GRAY},
        {'text': '  風險提醒', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  • 盤好時可能成為提款機（資金輪出）', 'size': 13, 'color': ACCENT_ORANGE},
        {'text': '  • 不要追高，等回檔再介入', 'size': 13, 'color': ACCENT_ORANGE},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(1.55), Inches(5.6), Inches(5.0),
                            leo_lines)

    # AI Cooling
    _add_rounded_rect(slide, Inches(6.8), Inches(1.4), Inches(6.2), Inches(5.3), WHITE)
    cool_lines = [
        {'text': 'AI 散熱（Thermal Solution）', 'size': 18, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  動能評級：★★★★☆', 'size': 15, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '', 'size': 3, 'color': DARK_GRAY},
        {'text': '  本週突破', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  • 散熱族群週二集體走出去（創新高）', 'size': 13, 'color': DARK_GRAY},
        {'text': '  • 規格確定下來 → 不確定性消除', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 3, 'color': DARK_GRAY},
        {'text': '  技術路線確認', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  • MCL（液冷微通道）太難 → 走 Cold Plate', 'size': 13, 'color': DARK_GRAY},
        {'text': '  • Stiffener + 2 片 Spreader 方案', 'size': 13, 'color': DARK_GRAY},
        {'text': '  • ASP 仍然很高，建策等龍頭受惠', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 3, 'color': DARK_GRAY},
        {'text': '  長期趨勢', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  • TDP 越來越高 → 靠近晶片端散熱是大瓶頸', 'size': 13, 'color': DARK_GRAY},
        {'text': '  • ASP + 訂單量都在拉高', 'size': 13, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  • QD 類（量大）人人能做，但高端散熱仍稀缺', 'size': 13, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(7.1), Inches(1.55), Inches(5.6), Inches(5.0),
                            cool_lines)


def slide_16_capex(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '科技巨頭資本支出軍備競賽', 'Google $1,800B | Amazon $2,000B — 2026 全面翻倍')
    _add_footer(slide, 16)

    img = charts['capex_war']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(10), Inches(3.8))

    _add_rounded_rect(slide, Inches(10.5), Inches(1.3), Inches(2.5), Inches(3.8), WHITE)
    note_lines = [
        {'text': '重點', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  調高 Capex', 'size': 12, 'color': CHARCOAL, 'bold': True},
        {'text': '  反而嚇壞', 'size': 12, 'color': CHARCOAL, 'bold': True},
        {'text': '  市場', 'size': 12, 'color': CHARCOAL, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  多頭時', 'size': 11, 'color': ACCENT_GREEN},
        {'text': '  = 利多', 'size': 11, 'color': ACCENT_GREEN},
        {'text': '', 'size': 3, 'color': DARK_GRAY},
        {'text': '  空頭時', 'size': 11, 'color': ACCENT_RED},
        {'text': '  = 燒錢', 'size': 11, 'color': ACCENT_RED},
    ]
    _add_multiline_text_box(slide, Inches(10.6), Inches(1.4), Inches(2.3), Inches(3.6),
                            note_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.4), Inches(12.7), Inches(1.4), WHITE)
    cap_bottom = [
        {'text': '本週 Capex 觀察', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  Google Q4 營收 $972B（超預期）+ 雲端 $177B（+48% YoY）→ 最強財報但 Capex $1,800B 嚇市場', 'size': 13, 'color': DARK_GRAY},
        {'text': '  Amazon 宣布 $2,000B Capex（比預期多 1/3）→ 盤後跌 10%，市場投反對票', 'size': 13, 'color': DARK_GRAY},
        {'text': '  核心判斷：AI 軍備競賽結束前誰能保住護城河才是關鍵，短期 Capex 恐慌是情緒性拋售', 'size': 13, 'color': ACCENT_GREEN, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(1.2),
                            cap_bottom)


def slide_17_part_five(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '05', 'PART FIVE.', '投資策略 — 本週佈局建議',
                         '推薦方向 | 資產配置 | 風險管理 | 下週觀察')


def slide_18_recommendations(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '本週推薦佈局方向', '美股 + 台股 — 族群性投資策略')
    _add_footer(slide, 18)

    img = charts['sector_radar']
    slide.shapes.add_picture(img, Inches(0.2), Inches(1.3), Inches(5.2), Inches(5.2))

    _add_rounded_rect(slide, Inches(5.6), Inches(1.3), Inches(7.4), Inches(2.3), WHITE)
    us_lines = [
        {'text': '美股推薦方向', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  AI 硬體（逢低佈局）：NVDA、AVGO、TSM、AMD', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  → 乖離超跌，PMI 上行支撐算力需求', 'size': 12, 'color': MID_GRAY},
        {'text': '  科技龍頭（長期持有）：GOOGL、AAPL、AMZN', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  → Google 最強勢，Apple 防禦佳', 'size': 12, 'color': MID_GRAY},
        {'text': '  傳產價值（趨勢正確）：KO、WMT、BRK.B', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  → 資金輪動受益，持續創歷史新高', 'size': 12, 'color': MID_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(5.8), Inches(1.4), Inches(7), Inches(2.1),
                            us_lines)

    _add_rounded_rect(slide, Inches(5.6), Inches(3.8), Inches(7.4), Inches(3.0), WHITE)
    tw_lines = [
        {'text': '台股推薦族群', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  低軌衛星（本週最強）：華通、耀華、同步基板', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  → 訂單確認，避風港效應，但盤好時注意獲利了結', 'size': 12, 'color': MID_GRAY},
        {'text': '  AI 散熱（突破確認中）：建策、奇鋐、超眾', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  → 規格確定（Cold Plate），ASP 拉高，觀察持續性', 'size': 12, 'color': MID_GRAY},
        {'text': '  記憶體（缺貨行情）：相關概念股', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  → Apple 被漲 70-80%，記憶體短缺打擊但也推升價格', 'size': 12, 'color': MID_GRAY},
        {'text': '  能源電力（中長期）：核能概念股', 'size': 14, 'color': CHARCOAL, 'bold': True},
        {'text': '  → AI 時代的「新石油」，3-5 年大趨勢', 'size': 12, 'color': MID_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(5.8), Inches(3.9), Inches(7), Inches(2.8),
                            tw_lines)


def slide_19_allocation(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '資產配置與風險管理', '穩健型 vs 積極型 — 根據投資屬性選擇')
    _add_footer(slide, 19)

    img = charts['allocation']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(8.5), Inches(3.5))

    _add_rounded_rect(slide, Inches(9.0), Inches(1.3), Inches(4.0), Inches(3.5), WHITE)
    risk_lines = [
        {'text': '風險管理要點', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  1. 波動放大 → 降槓桿', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  大 A 出現就先降部位', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  2. 看族群不看個股', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  用族群性角度判斷資金流向', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  3. 創新高的才追蹤', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  反彈不追，創新高才有持續性', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  4. 7:3 法則', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  70% 指數 + 30% 選股', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  5. 賣掉的就是非洲', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  減碼的不追回，等下個符合條件的', 'size': 12, 'color': MID_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(9.2), Inches(1.4), Inches(3.6), Inches(3.3),
                            risk_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), CORP_BLUE)
    summary = [
        {'text': '本週核心策略', 'size': 18, 'color': WHITE, 'bold': True},
        {'text': '', 'size': 5, 'color': WHITE},
        {'text': '  左側從來不悲觀，右側從來不樂觀 — 乖離超跌 + 情緒未失控 = 左側佈局機會浮現', 'size': 15, 'color': WHITE},
        {'text': '  保持本業收入，源源不絕的收入才有彎腰撿鑽石的空間', 'size': 15, 'color': SKY_BLUE},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), summary)


def slide_20_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CORP_BLUE)
    _add_circle(slide, Inches(10.5), Inches(-0.5), Inches(3), CORP_BLUE_DARK)
    _add_circle(slide, Inches(-0.8), Inches(5.5), Inches(2), CORP_BLUE_DARK)

    _add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
                  '本週總結 & 下週觀察重點', font_size=32, font_color=WHITE,
                  bold=True, alignment=PP_ALIGN.CENTER)
    _add_bg_rect(slide, Inches(5), Inches(1.3), Inches(3.3), Inches(0.04), WHITE)

    _add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.2), WHITE)
    summary_lines = [
        {'text': '本週五大重點回顧', 'size': 18, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  1. 川普推文救市：道瓊 +1,207 點創歷史新高 50,115，但科技股仍承壓（費半 -3.12%）', 'size': 14, 'color': DARK_GRAY},
        {'text': '  2. 風險資產退潮：比特幣 -30%、白銀 -17%、軟體股 -20% — 全面去槓桿而非單一利空', 'size': 14, 'color': DARK_GRAY},
        {'text': '  3. 日本大選變局：高市 352 席完全執政 → 國防擴張 + 半導體投資 + 消費稅減免', 'size': 14, 'color': DARK_GRAY},
        {'text': '  4. PMI 重返擴張：ISM 52.6 突破 3 年庫存調整，台灣 PMI 57 → 製造業上行循環確認', 'size': 14, 'color': DARK_GRAY},
        {'text': '  5. 產業爆發訊號：低軌衛星避風港 + 散熱族群突破 + 記憶體缺貨 → 輪動加速', 'size': 14, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(3.0),
                            summary_lines)

    _add_rounded_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.8), WHITE)
    next_lines = [
        {'text': '下週觀察重點', 'size': 18, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  美國 1 月非農就業數據（延至週三公佈）— 值缺數持續下滑的效果', 'size': 14, 'color': CHARCOAL},
        {'text': '  日本大選後政策落地速度 — 國防支出 & 消費稅減免時程', 'size': 14, 'color': CHARCOAL},
        {'text': '  科技股反彈力道是否持續 — 川普效應能撐多久？', 'size': 14, 'color': CHARCOAL},
        {'text': '  比特幣能否守住川普防線（$60K）— 200 週均線攻防戰', 'size': 14, 'color': CHARCOAL},
    ]
    _add_multiline_text_box(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.6),
                            next_lines)

    _add_text_box(slide, Inches(2), Inches(7.0), Inches(9.3), Inches(0.4),
                  '資料來源：財金號角（早晨財經速解讀 2/5, 2/6, 2/9）‧ 股癌投資 EP633 | 財金週報 W6',
                  font_size=11, font_color=SKY_BLUE, alignment=PP_ALIGN.CENTER)


# ============================================================
# Main Build
# ============================================================

def build_presentation():
    setup_matplotlib()
    print('Generating charts...')
    charts = gen_charts()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print('Building slides...')

    slide_01_cover(prs)
    print('  [1/20] Cover')

    slide_02_contents(prs)
    print('  [2/20] Contents')

    slide_03_part_one(prs)
    print('  [3/20] Part 1 divider')

    slide_04_us_weekly(prs, charts)
    print('  [4/20] US weekly indices')

    slide_05_asia(prs, charts)
    print('  [5/20] Asia markets')

    slide_06_part_two(prs)
    print('  [6/20] Part 2 divider')

    slide_07_trump_china(prs)
    print('  [7/20] Trump & China')

    slide_08_japan(prs, charts)
    print('  [8/20] Japan election')

    slide_09_bitcoin(prs, charts)
    print('  [9/20] Bitcoin crash')

    slide_10_part_three(prs)
    print('  [10/20] Part 3 divider')

    slide_11_pmi(prs, charts)
    print('  [11/20] ISM PMI')

    slide_12_value_growth(prs, charts)
    print('  [12/20] Value vs Growth')

    slide_13_ai_jobs(prs)
    print('  [13/20] AI & Jobs')

    slide_14_part_four(prs)
    print('  [14/20] Part 4 divider')

    slide_15_sectors(prs, charts)
    print('  [15/20] Sectors: LEO & Cooling')

    slide_16_capex(prs, charts)
    print('  [16/20] Tech Capex War')

    slide_17_part_five(prs)
    print('  [17/20] Part 5 divider')

    slide_18_recommendations(prs, charts)
    print('  [18/20] Recommendations')

    slide_19_allocation(prs, charts)
    print('  [19/20] Allocation')

    slide_20_summary(prs)
    print('  [20/20] Summary')

    prs.save(OUTPUT_PPTX)
    print(f'\nPresentation saved to: {OUTPUT_PPTX}')
    return OUTPUT_PPTX


if __name__ == '__main__':
    build_presentation()
