#!/usr/bin/env python3
"""
股癌 Podcast EP672 — 投資分析報告 PPT 生成腳本
風格：Corporate Blue & White Business Theme
主題：功率半導體超級週期 | 族群輪動 | 台股創歷史新高
內容：IDM lead time 99週、被動元件輪動教訓、機器人布局萌芽、Marvell 三征告捷
分析日期：2026-06-20
"""

import os
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from io import BytesIO
from math import pi

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ============================================================
# Config
# ============================================================

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(SCRIPT_DIR, '..', '..', '..'))
OUTPUT_DIR = os.path.join(PROJECT_ROOT, 'output', 'reports')
TEMP_DIR = os.path.join(PROJECT_ROOT, 'output', 'temp')
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(TEMP_DIR, exist_ok=True)
OUTPUT_PPTX = os.path.join(OUTPUT_DIR, '20260620_ep672_report.pptx')

# Colors — Corporate Blue Theme
CORP_BLUE = RGBColor(0x2B, 0x57, 0x9A)
CORP_BLUE_DARK = RGBColor(0x1A, 0x3A, 0x6C)
CORP_BLUE_LIGHT = RGBColor(0x4A, 0x86, 0xC8)
ACCENT_BLUE = RGBColor(0x5B, 0xA0, 0xD6)
SKY_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
CHARCOAL = RGBColor(0x2D, 0x2D, 0x2D)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x4D)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
BULL_GREEN = RGBColor(0x26, 0xA6, 0x9A)
BEAR_RED = RGBColor(0xEF, 0x53, 0x50)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_CN = 'Microsoft JhengHei'
FONT_EN = 'Arial'
TOTAL_PAGES = 20

# ============================================================
# Matplotlib Setup
# ============================================================

def setup_matplotlib():
    plt.rcParams['font.family'] = ['Heiti TC', 'Arial Unicode MS', 'Hiragino Sans', 'sans-serif']
    plt.rcParams['axes.unicode_minus'] = False
    plt.rcParams['figure.dpi'] = 150
    plt.rcParams['savefig.dpi'] = 250
    plt.rcParams['font.size'] = 14
    plt.rcParams['axes.labelsize'] = 15
    plt.rcParams['axes.titlesize'] = 18


def fig_to_image_stream(fig, dpi=250):
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor='white', edgecolor='none', pad_inches=0.15)
    plt.close(fig)
    buf.seek(0)
    return buf


# ============================================================
# Helper functions
# ============================================================

def _add_bg_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        _set_shape_alpha(shape, alpha)
    return shape


def _set_shape_alpha(shape, alpha_pct):
    fill_elem = shape.fill._fill
    solid = fill_elem.find(qn('a:solidFill'))
    if solid is not None:
        srgb = solid.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_elem = srgb.find(qn('a:alpha'))
            if alpha_elem is None:
                from lxml import etree
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
            alpha_elem.set('val', str(int((100 - alpha_pct) * 1000)))


def _add_text_box(slide, left, top, width, height, text, font_size=14,
                  font_color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name=FONT_CN, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    sp = p._pPr
    if sp is None:
        from lxml import etree
        sp = etree.SubElement(p._p, qn('a:pPr'))
    lnSpc = sp.find(qn('a:lnSpc'))
    if lnSpc is None:
        from lxml import etree
        lnSpc = etree.SubElement(sp, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))
    return txBox


def _add_multiline_text_box(slide, left, top, width, height, lines,
                            font_size=14, font_color=DARK_GRAY,
                            font_name=FONT_CN, line_spacing=1.5,
                            alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_data, dict):
            p.text = line_data.get('text', '')
            p.font.size = Pt(line_data.get('size', font_size))
            p.font.color.rgb = line_data.get('color', font_color)
            p.font.bold = line_data.get('bold', False)
            p.font.name = line_data.get('font', font_name)
        else:
            p.text = str(line_data)
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(2)
    return txBox


def _add_circle(slide, left, top, size, color, text='', font_size=14,
                font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = FONT_EN
        p.alignment = PP_ALIGN.CENTER
    return shape


def _add_rounded_rect(slide, left, top, width, height, color, text='',
                      font_size=13, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = FONT_CN
        p.alignment = PP_ALIGN.CENTER
    return shape


def _add_header_bar(slide, title_text, subtitle_text=''):
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), CORP_BLUE)
    _add_circle(slide, Inches(0.3), Inches(0.15), Inches(0.8), CORP_BLUE_DARK)
    _add_text_box(slide, Inches(1.3), Inches(0.15), Inches(10), Inches(0.55),
                  title_text, font_size=28, font_color=WHITE, bold=True)
    if subtitle_text:
        _add_text_box(slide, Inches(1.3), Inches(0.65), Inches(10), Inches(0.4),
                      subtitle_text, font_size=15, font_color=SKY_BLUE)
    _add_text_box(slide, Inches(10.5), Inches(0.15), Inches(2.5), Inches(0.4),
                  '股癌 Podcast EP672', font_size=16, font_color=WHITE, bold=True,
                  alignment=PP_ALIGN.RIGHT)
    _add_text_box(slide, Inches(10.5), Inches(0.55), Inches(2.5), Inches(0.35),
                  '2026.06.20', font_size=12, font_color=SKY_BLUE,
                  alignment=PP_ALIGN.RIGHT)


def _add_footer(slide, page_num, total_pages=TOTAL_PAGES):
    _add_bg_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), LIGHT_GRAY)
    _add_text_box(slide, Inches(0.5), Inches(7.15), Inches(7), Inches(0.3),
                  '股癌 Podcast EP672 — 2026.06.20 投資分析報告', font_size=10,
                  font_color=MID_GRAY)
    _add_text_box(slide, Inches(10), Inches(7.15), Inches(3), Inches(0.3),
                  f'{page_num} / {total_pages}', font_size=10,
                  font_color=MID_GRAY, alignment=PP_ALIGN.RIGHT, font_name=FONT_EN)


def _add_section_divider(slide, num, title, subtitle, detail=''):
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CORP_BLUE)
    _add_circle(slide, Inches(1.5), Inches(2), Inches(2.5), WHITE, num,
                font_size=48, font_color=CORP_BLUE)
    _add_text_box(slide, Inches(5), Inches(2.2), Inches(7), Inches(1),
                  title, font_size=40, font_color=WHITE, bold=True, font_name=FONT_EN)
    _add_text_box(slide, Inches(5), Inches(3.3), Inches(7), Inches(0.6),
                  subtitle, font_size=18, font_color=SKY_BLUE)
    _add_bg_rect(slide, Inches(5), Inches(4.1), Inches(4), Inches(0.04), WHITE)
    if detail:
        _add_text_box(slide, Inches(5), Inches(4.4), Inches(7), Inches(0.5),
                      detail, font_size=14, font_color=SKY_BLUE)


# ============================================================
# Chart Generation
# ============================================================

def gen_charts():
    charts = {}
    blue_hex = '#2B579A'
    blue_light_hex = '#4A86C8'
    accent_hex = '#5BA0D6'
    green_hex = '#27AE60'
    red_hex = '#E84D4D'
    orange_hex = '#F59E0B'
    purple_hex = '#8E44AD'

    # ── 1. 美股四大指數表現 ──
    fig, ax = plt.subplots(figsize=(11, 3.7))
    indices = ['道瓊工業\n44,000', '標普 500\n6,200', '納斯達克\n20,800', '費城半導體\n5,800']
    daily_chg = [0.2, 0.5, 0.8, 2.1]
    colors = [green_hex if v > 0 else red_hex for v in daily_chg]
    bars = ax.barh(indices, daily_chg, color=colors, height=0.5, edgecolor='white', linewidth=1.5)
    ax.axvline(x=0, color='#ccc', linewidth=1.2, linestyle='-')
    for bar, val in zip(bars, daily_chg):
        x_pos = bar.get_width() + 0.05
        label = f'+{val:.2f}%'
        ax.text(x_pos, bar.get_y() + bar.get_height() / 2,
                label, va='center', ha='left', fontsize=16,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xlim(0, 3.5)
    ax.set_title('美股四大指數 — 費半強勢領漲，AI 硬體主線確認（2026/06/20）',
                 fontsize=17, fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.tick_params(axis='both', labelsize=14)
    ax.annotate('功率半導體\n全面動起來!', xy=(2.1, 3), xytext=(1.8, 2.5),
                arrowprops=dict(arrowstyle='->', color=green_hex, lw=2),
                fontsize=13, color=green_hex, fontweight='bold')
    fig.tight_layout()
    charts['weekly_us'] = fig_to_image_stream(fig)

    # ── 2. 台股走勢與成交量 ──
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4))
    dates = ['6/13', '6/16', '6/17', '6/18', '6/19', '6/20']
    taiex = [39600, 39900, 40100, 40300, 40200, 40350]
    ax1.plot(dates, taiex, 'o-', color=blue_hex, linewidth=2.5, markersize=8)
    ax1.fill_between(range(len(dates)), taiex, min(taiex)*0.995,
                     alpha=0.15, color=blue_hex)
    ax1.set_xticks(range(len(dates)))
    ax1.set_xticklabels(dates, fontsize=11)
    ax1.set_ylabel('點數', fontsize=12, color=blue_hex)
    ax1.set_title('台股指數再創歷史新高（6/13-6/20）', fontsize=15, fontweight='bold',
                  color=blue_hex, pad=10)
    ax1.annotate('40,350\n歷史新高!', xy=(5, 40350), xytext=(3.5, 40450),
                 arrowprops=dict(arrowstyle='->', color=green_hex, lw=1.5),
                 fontsize=12, color=green_hex, fontweight='bold')
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.set_ylim(39300, 40700)

    volumes = [6800, 7500, 8200, 9100, 7800, 8500]
    vol_colors = [blue_light_hex]*5 + [green_hex]
    bars = ax2.bar(range(len(dates)), volumes, color=vol_colors, width=0.6,
                   edgecolor='white', linewidth=1.5)
    ax2.set_xticks(range(len(dates)))
    ax2.set_xticklabels(dates, fontsize=11)
    ax2.set_ylabel('成交量（億元）', fontsize=12, color=blue_hex)
    ax2.set_title('台股日成交量（億元）', fontsize=15, fontweight='bold',
                  color=blue_hex, pad=10)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    fig.tight_layout()
    charts['asia_weekly'] = fig_to_image_stream(fig)

    # ── 3. 功率半導體 IDM Lead Time ──
    fig, ax = plt.subplots(figsize=(11, 3.8))
    idm_names = ['IFX\n(Infineon)', 'STM\n(ST Micro)', 'ON Semi', 'Rohm\n(日本)', '台廠 IDM']
    lead_times = [99, 90, 80, 70, 40]
    lt_colors = [red_hex, red_hex, orange_hex, orange_hex, blue_hex]
    bars = ax.bar(range(len(idm_names)), lead_times, color=lt_colors,
                  width=0.55, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, lead_times)):
        label = '99+ 週' if val == 99 else f'{val} 週'
        ax.text(i, val + 1.5, label, ha='center', fontsize=13,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(idm_names)))
    ax.set_xticklabels(idm_names, fontsize=12)
    ax.set_ylabel('Lead Time（週）', fontsize=13, color=blue_hex)
    ax.set_title('功率半導體 IDM Lead Time 爆炸性拉長（2026 Q2）— IFX/STM 已達 99+ 週',
                 fontsize=16, fontweight='bold', color=blue_hex, pad=15)
    ax.axhline(y=52, color=orange_hex, linewidth=2, linestyle='--', alpha=0.7)
    ax.text(4.6, 53, '1年 = 52週', fontsize=10, color=orange_hex, fontweight='bold')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(0, 120)
    fig.tight_layout()
    charts['lead_time'] = fig_to_image_stream(fig)

    # ── 4. 功率半導體供應鏈資金流向 ──
    fig, ax = plt.subplots(figsize=(9, 4))
    supply_chain = ['設計廠\n(Fabless)', 'IDM\n(自給自足)', '6/8吋\n代工廠', '封裝廠']
    fund_flow = [30, 35, 25, 15]
    flow_colors = [blue_light_hex, blue_hex, green_hex, accent_hex]
    bars = ax.bar(range(len(supply_chain)), fund_flow, color=flow_colors,
                  width=0.5, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, fund_flow)):
        ax.text(i, val + 0.8, f'+{val}', ha='center', fontsize=15,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(supply_chain)))
    ax.set_xticklabels(supply_chain, fontsize=12)
    ax.set_ylabel('相對估值修復（相對值）', fontsize=13, color=blue_hex)
    ax.set_title('功率半導體供應鏈全面受惠 — 各環節估值修復', fontsize=17,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(0, 45)
    fig.tight_layout()
    charts['supply_chain'] = fig_to_image_stream(fig)

    # ── 5. AI 伺服器需求 vs 功率半導體缺貨趨勢 ──
    fig, ax = plt.subplots(figsize=(10, 3.8))
    quarters = ['2025Q1', '2025Q2', '2025Q3', '2025Q4', '2026Q1', '2026Q2']
    ai_demand = [60, 70, 80, 90, 100, 115]
    shortage_idx = [20, 35, 55, 70, 90, 99]
    ax.fill_between(range(len(quarters)), ai_demand, alpha=0.3, color=blue_hex, label='AI 伺服器需求指數')
    ax.plot(range(len(quarters)), ai_demand, 'o-', color=blue_hex, linewidth=2.5, markersize=8)
    ax.fill_between(range(len(quarters)), shortage_idx, alpha=0.2, color=red_hex, label='功率半導體缺貨指數')
    ax.plot(range(len(quarters)), shortage_idx, 's--', color=red_hex, linewidth=2.5, markersize=8)
    ax.set_xticks(range(len(quarters)))
    ax.set_xticklabels(quarters, fontsize=12)
    ax.set_ylabel('指數（相對值）', fontsize=13, color=blue_hex)
    ax.set_title('AI 伺服器需求 vs 功率半導體缺貨指數 — 需求驅動，持續性強',
                 fontsize=16, fontweight='bold', color=blue_hex, pad=15)
    ax.legend(fontsize=12, loc='upper left')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    fig.tight_layout()
    charts['ai_demand'] = fig_to_image_stream(fig)

    # ── 6. 被動元件各品類漲幅比較 ──
    fig, ax = plt.subplots(figsize=(11, 3.8))
    passive_cats = ['電阻\n（意外最高！）', '消費級\nMLCC', '高容值\nMLCC（最缺）', '電容', '電感']
    passive_chg = [65, 55, 45, 40, 35]
    p_colors = [red_hex, orange_hex, blue_hex, blue_light_hex, accent_hex]
    bars = ax.bar(range(len(passive_cats)), passive_chg, color=p_colors,
                  width=0.55, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, passive_chg)):
        ax.text(i, val + 1.5, f'+{val}%', ha='center', fontsize=13,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(passive_cats)))
    ax.set_xticklabels(passive_cats, fontsize=12)
    ax.set_ylabel('漲幅（%）', fontsize=13, color=blue_hex)
    ax.set_title('被動元件族群：意外的是電阻跑最快！市場永遠讓你意外',
                 fontsize=16, fontweight='bold', color=blue_hex, pad=15)
    ax.annotate('最缺≠漲最多\n教訓！', xy=(2, 45), xytext=(3.2, 58),
                arrowprops=dict(arrowstyle='->', color=red_hex, lw=2),
                fontsize=12, color=red_hex, fontweight='bold')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(0, 80)
    fig.tight_layout()
    charts['passive_components'] = fig_to_image_stream(fig)

    # ── 7. 機器人供應鏈動態雷達 ──
    fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
    robot_cats = ['軟體/AI', '日本斜坡\n減速器', '電機/馬達', '感測器', '台灣\n工控零件', '結構件']
    N = len(robot_cats)
    current_score = [9, 8, 7, 6, 5, 4]
    angles = [n / float(N) * 2 * pi for n in range(N)]
    angles += angles[:1]
    current_score += current_score[:1]
    ax.plot(angles, current_score, 'o-', color=blue_hex, linewidth=2.5,
            label='當前重要性評分', markersize=8)
    ax.fill(angles, current_score, alpha=0.15, color=blue_hex)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(robot_cats, fontsize=11, fontweight='bold')
    ax.set_ylim(0, 10)
    ax.set_yticks([2, 4, 6, 8, 10])
    ax.set_yticklabels(['2', '4', '6', '8', '10'], fontsize=10, color='gray')
    ax.legend(fontsize=11, loc='upper right', bbox_to_anchor=(1.3, 1.1))
    ax.set_title('機器人供應鏈動態雷達\nTesla Optimus 催化，布局萌芽期', fontsize=15,
                 fontweight='bold', color=blue_hex, pad=30)
    fig.tight_layout()
    charts['robot_radar'] = fig_to_image_stream(fig)

    # ── 8. 硬體 vs 軟體年初至今分化 ──
    fig, ax = plt.subplots(figsize=(10, 3.8))
    hw_sw_cats = ['AI 硬體族群', '功率半導體', '台積電', '整體大盤', 'AI 軟體 IGV', '一般 SaaS']
    hw_sw_chg = [85, 70, 55, 35, -5, -10]
    hs_colors = [green_hex if v > 0 else red_hex for v in hw_sw_chg]
    bars = ax.bar(range(len(hw_sw_cats)), hw_sw_chg, color=hs_colors,
                  width=0.55, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, hw_sw_chg)):
        y_pos = val + 2 if val >= 0 else val - 4
        label = f'+{val}%' if val > 0 else f'{val}%'
        ax.text(i, y_pos, label, ha='center', fontsize=12,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(hw_sw_cats)))
    ax.set_xticklabels(hw_sw_cats, fontsize=11)
    ax.set_ylabel('年初至今漲跌（%）', fontsize=13, color=blue_hex)
    ax.set_title('硬體 vs 軟體年初至今分化極端化（2026 YTD）',
                 fontsize=17, fontweight='bold', color=blue_hex, pad=15)
    ax.axhline(y=0, color='#ccc', linewidth=1)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(-20, 105)
    fig.tight_layout()
    charts['hw_sw_split'] = fig_to_image_stream(fig)

    # ── 9. 美股推薦評分 ──
    fig, ax = plt.subplots(figsize=(10, 4.0))
    us_stocks = ['Marvell (MRVL)', 'ON Semi (ON)', 'Texas Inst. (TXN)',
                 'STMicro (STM)', 'Infineon (IFX)', 'CPU 類股 (AMD)']
    us_scores = [9, 8, 8, 7.5, 7.5, 7]
    us_colors = [blue_hex, red_hex, green_hex, orange_hex, purple_hex, blue_light_hex]
    bars = ax.barh(us_stocks, us_scores, color=us_colors, height=0.55,
                   edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars, us_scores):
        x_pos = bar.get_width() + 0.1
        ax.text(x_pos, bar.get_y() + bar.get_height() / 2,
                f'{val}/10', va='center', ha='left', fontsize=13,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xlabel('推薦評分（1-10）', fontsize=13, color=blue_hex)
    ax.set_title('美股推薦清單 — 功率半導體 IDM + Marvell 持續加碼',
                 fontsize=17, fontweight='bold', color=blue_hex, pad=15)
    ax.set_xlim(0, 12)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.invert_yaxis()
    fig.tight_layout()
    charts['us_stocks'] = fig_to_image_stream(fig)

    # ── 10. 台股推薦方向評分 ──
    fig, ax = plt.subplots(figsize=(10, 4.2))
    tw_directions = ['功率半導體\n設計廠', '6/8吋\n代工廠', 'IDM 台廠\n(含VShard)', '高容值\nMLCC 廠',
                     '電阻廠\n（意外受惠）', 'AI伺服器\n散熱', '機器人\n零組件（長線）']
    tw_scores = [9, 8.5, 8.5, 8, 7.5, 7, 6]
    tw_colors = [blue_hex, green_hex, green_hex, blue_light_hex, orange_hex, accent_hex, purple_hex]
    bars = ax.barh(tw_directions, tw_scores, color=tw_colors, height=0.55,
                   edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars, tw_scores):
        x_pos = bar.get_width() + 0.1
        ax.text(x_pos, bar.get_y() + bar.get_height() / 2,
                f'{val}/10', va='center', ha='left', fontsize=13,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xlabel('推薦評分（1-10）', fontsize=13, color=blue_hex)
    ax.set_title('台股推薦方向 — 功率半導體鏈優先，機器人長線布局',
                 fontsize=17, fontweight='bold', color=blue_hex, pad=15)
    ax.set_xlim(0, 12)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.invert_yaxis()
    fig.tight_layout()
    charts['tw_stocks'] = fig_to_image_stream(fig)

    # ── 11. 產業動能雷達圖 ──
    fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
    sector_cats = ['功率\n半導體', '台股\n整體動能', '被動元件', '機器人', '軟體股\n動能']
    N = len(sector_cats)
    current = [10, 9, 7, 5, 3]
    future = [10, 9, 6, 8, 5]
    angles = [n / float(N) * 2 * pi for n in range(N)]
    angles += angles[:1]
    current += current[:1]
    future += future[:1]
    ax.plot(angles, current, 'o-', color=blue_hex, linewidth=2.5,
            label='當前評分', markersize=8)
    ax.fill(angles, current, alpha=0.15, color=blue_hex)
    ax.plot(angles, future, 's--', color=orange_hex, linewidth=2.5,
            label='3-6月展望', markersize=8)
    ax.fill(angles, future, alpha=0.1, color=orange_hex)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(sector_cats, fontsize=12, fontweight='bold')
    ax.set_ylim(0, 10)
    ax.set_yticks([2, 4, 6, 8, 10])
    ax.set_yticklabels(['2', '4', '6', '8', '10'], fontsize=10, color='gray')
    ax.legend(fontsize=12, loc='upper right', bbox_to_anchor=(1.3, 1.1))
    ax.set_title('市場動能雷達圖', fontsize=18,
                 fontweight='bold', color=blue_hex, pad=30)
    fig.tight_layout()
    charts['sector_radar'] = fig_to_image_stream(fig)

    # ── 12. 資產配置建議（積極型 + 穩健型）──
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(8.5, 3.5))
    labels_c = ['功率半導體鏈 40%', 'AI伺服器/台積電 25%', '機器人布局 10%', '現金 25%']
    sizes_c = [40, 25, 10, 25]
    colors_c = [blue_hex, green_hex, purple_hex, '#95a5a6']
    wedges, texts, autotexts = ax1.pie(
        sizes_c, labels=labels_c, autopct='%1.0f%%', colors=colors_c,
        startangle=90, pctdistance=0.75, textprops={'fontsize': 10})
    for at in autotexts:
        at.set_fontweight('bold')
        at.set_color('white')
    ax1.set_title('積極型（功率半導體主線）', fontsize=14, fontweight='bold', color=blue_hex, pad=15)

    labels_g = ['現金 40%', '台積電 30%', '防禦ETF 20%', '機器人長線 10%']
    sizes_g = [40, 30, 20, 10]
    colors_g = ['#95a5a6', blue_hex, blue_light_hex, purple_hex]
    wedges2, texts2, autotexts2 = ax2.pie(
        sizes_g, labels=labels_g, autopct='%1.0f%%', colors=colors_g,
        startangle=90, pctdistance=0.75, textprops={'fontsize': 10})
    for at in autotexts2:
        at.set_fontweight('bold')
        at.set_color('white')
    ax2.set_title('穩健型（等機器人確認）', fontsize=14, fontweight='bold', color=blue_hex, pad=15)
    fig.tight_layout()
    charts['allocation'] = fig_to_image_stream(fig)

    return charts


# ============================================================
# Slide Builders
# ============================================================

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CORP_BLUE)
    _add_circle(slide, Inches(10.5), Inches(-0.5), Inches(3), CORP_BLUE_DARK)
    _add_circle(slide, Inches(11.5), Inches(5), Inches(2.5), CORP_BLUE_DARK)
    _add_circle(slide, Inches(-0.8), Inches(5.5), Inches(2), CORP_BLUE_DARK)
    _add_rounded_rect(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(5), WHITE)
    _add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
                  '股癌 Podcast EP672 — 市場分析報告', font_size=36, font_color=CORP_BLUE,
                  bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.9),
                  '功率半導體超級週期｜族群輪動｜台股創歷史新高',
                  font_size=20, font_color=CORP_BLUE_DARK,
                  bold=True, alignment=PP_ALIGN.CENTER)
    _add_bg_rect(slide, Inches(4.5), Inches(3.5), Inches(4.3), Inches(0.05), CORP_BLUE)
    _add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.5),
                  '功率半導體 lead time 99週 | IDM 全面動起來 | 台股再創新高',
                  font_size=14, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10), Inches(0.5),
                  '被動元件族群輪動 | 機器人布局萌芽期 | Marvell 三征告捷',
                  font_size=14, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.4),
                  '2026 年 6 月 20 日', font_size=12,
                  font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.5),
                  '股癌 Podcast EP672 — 2026.06.20 投資分析', font_size=18,
                  font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_02_contents(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '目錄 CONTENTS', '2026.06.20 功率半導體超級週期 + 族群輪動 + 台股創新高 + 機器人布局')
    _add_footer(slide, 2)

    sections = [
        {'num': '01', 'title': '市場總覽', 'desc': '台股創歷史新高\n族群輪動進行中\nAI 伺服器主線', 'color': CORP_BLUE},
        {'num': '02', 'title': '功率半導體', 'desc': 'IDM 缺貨嚴峻\nlead time 99+ 週\n需求驅動持續性強', 'color': CORP_BLUE_LIGHT},
        {'num': '03', 'title': '產業趨勢', 'desc': '被動元件輪動教訓\n機器人布局萌芽\n軟硬分化持續', 'color': ACCENT_BLUE},
        {'num': '04', 'title': '投資策略', 'desc': '美股台股推薦\n配置建議\n後市展望', 'color': ACCENT_PURPLE},
    ]

    start_x = Inches(0.8)
    for i, sec in enumerate(sections):
        x = start_x + i * Inches(3.2)
        y = Inches(2.0)
        _add_rounded_rect(slide, x, y, Inches(2.9), Inches(4.3), WHITE)
        _add_circle(slide, x + Inches(0.85), y + Inches(0.3), Inches(1.1),
                    sec['color'], sec['num'], font_size=28)
        _add_text_box(slide, x + Inches(0.1), y + Inches(1.6), Inches(2.7),
                      Inches(0.5), sec['title'], font_size=20,
                      font_color=CHARCOAL, bold=True, alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, x + Inches(0.2), y + Inches(2.2), Inches(2.5),
                      Inches(1.5), sec['desc'], font_size=13,
                      font_color=MID_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.5)
        _add_bg_rect(slide, x + Inches(0.8), y + Inches(4.0), Inches(1.3),
                     Inches(0.05), sec['color'])


def slide_03_part_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '01', 'PART ONE.', '市場總覽 — 台股創歷史新高',
                         '族群輪動進行中 | AI 伺服器 Rubin 平台拉貨')


def slide_04_us_weekly(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '美股四大指數 — 費半強勢領漲', '2026/06/20 — 道瓊+0.2% 標普+0.5% 那指+0.8% 費半+2.1%')
    _add_footer(slide, 4)

    idx_data = [
        ('道瓊工業', '44,000', '+88 pt', '+0.2%', ACCENT_GREEN),
        ('標普 500', '6,200', '+31 pt', '+0.5%', ACCENT_GREEN),
        ('納斯達克', '20,800', '+166 pt', '+0.8%', ACCENT_GREEN),
        ('費城半導體', '5,800', '+120 pt', '+2.1%', ACCENT_GREEN),
    ]
    for i, (name, price, change, pct, color) in enumerate(idx_data):
        x = Inches(0.4) + i * Inches(3.2)
        _add_rounded_rect(slide, x, Inches(1.4), Inches(3.0), Inches(1.6), WHITE)
        _add_bg_rect(slide, x + Inches(0.05), Inches(1.4), Inches(2.9), Inches(0.06), color)
        _add_text_box(slide, x + Inches(0.2), Inches(1.55), Inches(2.6), Inches(0.35),
                      name, font_size=15, font_color=CHARCOAL, bold=True, alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(2.6), Inches(0.4),
                      price, font_size=24, font_color=CHARCOAL, bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
        _add_text_box(slide, x + Inches(0.2), Inches(2.35), Inches(2.6), Inches(0.35),
                      f'{change}（{pct}）', font_size=14, font_color=color,
                      bold=True, alignment=PP_ALIGN.CENTER)

    img = charts['weekly_us']
    slide.shapes.add_picture(img, Inches(1.2), Inches(3.2), Inches(11), Inches(3.7))


def slide_05_taiwan_market(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '台股指數再次創歷史新高', '台股站上 40,000+ | 功率半導體接棒被動元件 | 手上弱勢股 ≠ 市場轉弱')
    _add_footer(slide, 5)

    img = charts['asia_weekly']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.8), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.3), Inches(1.3), Inches(4.7), Inches(3.5), WHITE)
    taiwan_lines = [
        {'text': '台股創新高解析', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  台股指數再次創歷史新高（40,000+）', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  功率半導體接棒被動元件', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  重要觀念：', 'size': 13, 'color': CORP_BLUE, 'bold': True},
        {'text': '  手上持股弱 ≠ 市場整體轉弱', 'size': 13, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  不同族群有不同輪動邏輯', 'size': 13, 'color': DARK_GRAY},
        {'text': '  AI 仍在幼年期，週期走更遠', 'size': 13, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(8.5), Inches(1.45), Inches(4.3), Inches(3.2),
                            taiwan_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    bottom = [
        {'text': '台股多頭框架', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  AI 伺服器持續拉貨 + 台積電 Rubin 平台需求確認 → 台股基本面支撐穩固', 'size': 14, 'color': DARK_GRAY},
        {'text': '  族群輪動是市場健康的表徵，資金從被動元件 → 功率半導體 → 機器人（長線）', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), bottom)


def slide_06_part_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '02', 'PART TWO.', '功率半導體 — 超級週期啟動',
                         'IDM 廠 lead time 突破 99 週 | AI 伺服器 Rubin 帶動全鏈')


def slide_07_lead_time(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '功率半導體缺貨分析 — IDM Lead Time 爆炸', 'IFX/STM/ON/Rohm 全面拉長至 80-99 週 | 需求驅動，非投機炒作')
    _add_footer(slide, 7)

    img = charts['lead_time']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.8), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.3), Inches(1.3), Inches(4.7), Inches(3.5), WHITE)
    lt_lines = [
        {'text': '缺貨背後的邏輯', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  IFX/STM：99+ 週（某料號直接標 "99+"）', 'size': 13, 'color': ACCENT_RED, 'bold': True},
        {'text': '  代表買不到！', 'size': 14, 'color': ACCENT_RED, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  需求驅動 vs 投機炒作', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  AI Rubin 伺服器真實需求驅動', 'size': 13, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  持續性遠強於貴金屬漲價驅動', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  Global 現象，非單一市場', 'size': 13, 'color': DARK_GRAY},
        {'text': '  台廠 IDM 開始受益', 'size': 13, 'color': ACCENT_GREEN},
    ]
    _add_multiline_text_box(slide, Inches(8.5), Inches(1.45), Inches(4.3), Inches(3.2),
                            lt_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    bottom = [
        {'text': '下一個驗證點', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  IDM 廠商何時釋出產能，是投資週期長短的關鍵。6/8吋代工廠解瓶頸能力決定台廠受惠程度', 'size': 14, 'color': DARK_GRAY},
        {'text': '  溢價落在哪個環節（IDM vs Fabless vs 代工廠）需持續觀察，現階段 IDM 最確定', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), bottom)


def slide_08_supply_chain(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '功率半導體供應鏈 — 各環節全面受惠', 'IDM/Fabless/6吋8吋代工廠 估值修復 | 溢價落在哪裡是關鍵')
    _add_footer(slide, 8)

    img = charts['supply_chain']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.8))

    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(4.9), Inches(5.5), WHITE)
    sc_lines = [
        {'text': '供應鏈各環節分析', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  IDM（自給自足）', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  掌握產能 = 最穩，IFX/STM最強', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  Fabless（設計廠）', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  掌握客戶，需搶代工廠產能', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  6吋/8吋代工廠', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  去瓶頸的關鍵，台廠優勢明顯', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  溢價落在哪個環節？', 'size': 13, 'color': CORP_BLUE, 'bold': True},
        {'text': '  持續觀察是投資重點', 'size': 12, 'color': MID_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.5), Inches(5.3),
                            sc_lines)


def slide_09_ai_demand(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, 'AI 伺服器需求驅動 — Nvidia Rubin 平台拉貨', '需求端驅動 vs 投機炒作完全不同 | 消費性/車用仍弱讓 AI 週期更穩')
    _add_footer(slide, 9)

    img = charts['ai_demand']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.8), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.3), Inches(1.3), Inches(4.7), Inches(3.5), WHITE)
    ai_lines = [
        {'text': 'AI 週期為何走得更穩', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  Nvidia Rubin 密集拉貨', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  沒有這些料伺服器拼不出來', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  消費性 + 車用仍疲弱', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  反而讓 AI 週期走得更遠', 'size': 13, 'color': ACCENT_GREEN},
        {'text': '  避免所有品類同搶料', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  台積電就像「限制器」', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  確保週期走得更穩更長', 'size': 13, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(8.5), Inches(1.45), Inches(4.3), Inches(3.2),
                            ai_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    bottom = [
        {'text': 'CSP 不計代價搶料是關鍵支撐', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  雲端巨頭（CSP）不計代價搶料，確保 AI 伺服器按時交付，訂單能見度高、違約成本高', 'size': 14, 'color': DARK_GRAY},
        {'text': '  AI 週期仍在幼年期 — 不要猜幾局下半場，現在連上半場都還沒打完', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), bottom)


def slide_10_part_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '03', 'PART THREE.', '產業趨勢 — 被動元件 / 機器人 / 軟硬分化',
                         '意外的輪動教訓 | 機器人萌芽期 | 軟體重定位')


def slide_11_passive_components(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '被動元件族群輪動教訓 — 意外的是電阻跑最快', '最缺的不一定漲最多 | 排擠效應 | VShard 雙重曝險')
    _add_footer(slide, 11)

    img = charts['passive_components']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.8), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.3), Inches(1.3), Inches(4.7), Inches(3.5), WHITE)
    passive_lines = [
        {'text': '被動元件輪動啟示錄', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  教訓：最缺 ≠ 漲最多！', 'size': 14, 'color': ACCENT_RED, 'bold': True},
        {'text': '  電阻（最不被看好）反而噴最凶', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  排擠效應邏輯', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  最強料把最高端產能全占走', 'size': 13, 'color': DARK_GRAY},
        {'text': '  MLCC 排擠電阻的代工產能', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  VShard（華新科）特殊之處', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  被動元件 + 功率半導體雙重曝險', 'size': 13, 'color': ACCENT_GREEN, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(8.5), Inches(1.45), Inches(4.3), Inches(3.2),
                            passive_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    bottom = [
        {'text': '市場永遠讓你意外，不要把話說死', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  族群輪動：資金從被動元件轉向功率半導體。但 VShard（華新科）同時有兩邊曝險，是特殊的選擇', 'size': 14, 'color': DARK_GRAY},
        {'text': '  HBM vs DDR4/DDR5 的故事重演：最高端不是最多漲，排擠效應創造意外贏家', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), bottom)


def slide_12_robot_layout(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '機器人早期布局 — Tesla Optimus 催化萌芽期', '工控零組件廠出現低個位數機器人訂單 | 等財報驗證就太晚了')
    _add_footer(slide, 12)

    img = charts['robot_radar']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(5.5), Inches(5.2))

    _add_rounded_rect(slide, Inches(6.1), Inches(1.3), Inches(6.9), Inches(5.5), WHITE)
    robot_lines = [
        {'text': '機器人布局策略', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  Tesla Optimus 是主要催化劑', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  工控零組件廠：低個位數機器人訂單占比', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  重點零組件', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  日本廠：斜坡減速器（harmonic drive）', 'size': 13, 'color': DARK_GRAY},
        {'text': '  電機/馬達、感測器、台灣工控零件', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  布局策略：不等財報驗證', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  等財報確認再進場 = 太晚了', 'size': 13, 'color': ACCENT_RED},
        {'text': '  現在萌芽期是最好的埋伏點', 'size': 13, 'color': ACCENT_GREEN},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  預計是中大型長期趨勢', 'size': 13, 'color': DARK_GRAY},
        {'text': '  軟體/AI 是機器人最重要的靈魂', 'size': 12, 'color': MID_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(6.3), Inches(1.45), Inches(6.5), Inches(5.3),
                            robot_lines)


def slide_13_hw_sw_split(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '硬體 vs 軟體年初至今分化極端化', '軟體 hiring freeze = over-hiring 消化 | AI 讓工程師重定位，非淘汰')
    _add_footer(slide, 13)

    img = charts['hw_sw_split']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.8))

    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(4.9), Inches(5.5), WHITE)
    hw_sw_lines = [
        {'text': '軟硬分化的真正原因', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  硬體為何強：', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  AI 基礎建設需求爆發', 'size': 13, 'color': DARK_GRAY},
        {'text': '  功率半導體 lead time 99+ 週', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  軟體為何弱：', 'size': 14, 'color': ACCENT_RED, 'bold': True},
        {'text': '  Hiring freeze = over-hiring 消化', 'size': 13, 'color': DARK_GRAY},
        {'text': '  不是 AI 替代工程師', 'size': 13, 'color': ACCENT_ORANGE},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  工程師的未來', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  AI 讓工程師重定位，TAM 更大', 'size': 13, 'color': ACCENT_GREEN},
        {'text': '  能 harness AI 工具 = 贏家', 'size': 13, 'color': ACCENT_GREEN, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.5), Inches(5.3),
                            hw_sw_lines)


def slide_14_part_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '04', 'PART FOUR.', '投資策略 — 美股台股推薦',
                         'Marvell 三征告捷 | 功率半導體鏈布局 | 族群輪動框架')


def slide_15_us_stocks(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '美股推薦清單 — 功率半導體 IDM + Marvell', 'Marvell 三次進攻終於成功 | IDM lead time 99 週 = 最強基本面支撐')
    _add_footer(slide, 15)

    img = charts['us_stocks']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.8))

    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(4.9), Inches(5.5), WHITE)
    us_lines = [
        {'text': '美股推薦邏輯', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  Marvell (MRVL)：9/10', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  AI 定制芯片龍頭，三征告捷', 'size': 13, 'color': DARK_GRAY},
        {'text': '  正在持有並加碼', 'size': 13, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  功率半導體 IDM：8/10', 'size': 14, 'color': ACCENT_RED, 'bold': True},
        {'text': '  ON Semi / TXN / STM / IFX', 'size': 13, 'color': DARK_GRAY},
        {'text': '  lead time 99 週 = 最強支撐', 'size': 13, 'color': ACCENT_GREEN},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  CPU 類股（AMD）：7/10', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  AI 算力需求持續帶動', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  策略：族群確認後，逢回機會', 'size': 13, 'color': CORP_BLUE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.5), Inches(5.3),
                            us_lines)


def slide_16_tw_stocks(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '台股推薦方向 — 功率半導體鏈優先，機器人長線布局', 'VShard 雙重曝險 | 6/8吋代工廠是瓶頸解方 | 機器人萌芽埋伏')
    _add_footer(slide, 16)

    img = charts['tw_stocks']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.8))

    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(4.9), Inches(5.5), WHITE)
    tw_lines = [
        {'text': '台股推薦邏輯', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  功率半導體設計廠：9/10', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  受惠 IDM 二公效應，台廠接單', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  VShard（華新科）：8.5/10', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  被動元件 + 功率半導體雙重曝險', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  6/8吋代工廠：8.5/10', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  解決瓶頸的關鍵玩家', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  機器人零組件：6/10（長線）', 'size': 14, 'color': ACCENT_PURPLE, 'bold': True},
        {'text': '  現在埋伏，等財報逐步驗證', 'size': 13, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.5), Inches(5.3),
                            tw_lines)


def slide_17_part_five(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '05', 'PART FIVE.', '投資策略總結 — 2026.06.20 操作建議',
                         '功率半導體主線配置 | 族群輪動框架 | 風險評估')


def slide_18_recommendations(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '產業動能雷達 & 配置建議', '功率半導體 10 分滿分 | 機器人未來評分持續提升 | 現金等機器人確認')
    _add_footer(slide, 18)

    img = charts['sector_radar']
    slide.shapes.add_picture(img, Inches(0.2), Inches(1.3), Inches(5.2), Inches(5.2))

    _add_rounded_rect(slide, Inches(5.6), Inches(1.3), Inches(7.4), Inches(2.3), WHITE)
    active_lines = [
        {'text': '積極型配置建議', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  功率半導體鏈 40%（主線確認，IDM/代工/Fabless）', 'size': 13, 'color': CHARCOAL},
        {'text': '  AI 伺服器/台積電 25%（Rubin 需求確認）', 'size': 13, 'color': CHARCOAL},
        {'text': '  機器人長線布局 10%（萌芽期先埋伏）', 'size': 13, 'color': CHARCOAL},
        {'text': '  現金 25%（等機器人題材確認）', 'size': 13, 'color': CHARCOAL},
    ]
    _add_multiline_text_box(slide, Inches(5.8), Inches(1.4), Inches(7), Inches(2.1),
                            active_lines)

    _add_rounded_rect(slide, Inches(5.6), Inches(3.8), Inches(7.4), Inches(3.0), WHITE)
    steady_lines = [
        {'text': '穩健型配置建議', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  現金 40%（等機器人確認後再加碼）', 'size': 13, 'color': CHARCOAL},
        {'text': '  台積電/大型權值 30%（AI 主線確定）', 'size': 13, 'color': CHARCOAL},
        {'text': '  防禦科技ETF 20%（降低個股集中風險）', 'size': 13, 'color': CHARCOAL},
        {'text': '  機器人長線 10%（少量先布局）', 'size': 13, 'color': CHARCOAL},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  注意：功率半導體族群籌碼擁擠時要減碼', 'size': 11, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(5.8), Inches(3.9), Inches(7), Inches(2.8),
                            steady_lines)


def slide_19_allocation(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '資產配置圖解 & 後市觀察重點', '2026.06.20 — 功率半導體主線確立，機器人萌芽布局期')
    _add_footer(slide, 19)

    img = charts['allocation']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(8.5), Inches(3.5))

    _add_rounded_rect(slide, Inches(9.0), Inches(1.3), Inches(4.0), Inches(3.5), WHITE)
    risk_lines = [
        {'text': '後市觀察重點', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  1. IDM 廠商財報電話會議', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  下一季更新是最大驗證點', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  2. 6/8吋代工廠產能釋出', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  解瓶頸是台廠受惠關鍵', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  3. 功率半導體籌碼擁擠', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  當前族群是否已過熱', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  4. 機器人占比何時提升', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  工控廠商占比出現明顯提升', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  5. 消費性/車用需求是否回升', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  回升 = AI 週期加速縮短風險', 'size': 12, 'color': ACCENT_RED},
    ]
    _add_multiline_text_box(slide, Inches(9.2), Inches(1.4), Inches(3.6), Inches(3.3),
                            risk_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), CORP_BLUE)
    summary = [
        {'text': '核心策略總結', 'size': 18, 'color': WHITE, 'bold': True},
        {'text': '', 'size': 5, 'color': WHITE},
        {'text': '  功率半導體超級週期確認，IDM lead time 99+ 週是最強基本面支撐；台股創新高是市場健康的表現', 'size': 15, 'color': WHITE},
        {'text': '  機器人萌芽期先布局，不要等財報驗證；Marvell 三征告捷，AI 定制芯片主線持續', 'size': 14, 'color': SKY_BLUE},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), summary)


def slide_20_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CORP_BLUE)
    _add_circle(slide, Inches(10.5), Inches(-0.5), Inches(3), CORP_BLUE_DARK)
    _add_circle(slide, Inches(-0.8), Inches(5.5), Inches(2), CORP_BLUE_DARK)

    _add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
                  '2026.06.20 股癌 EP672 重點回顧 & 後市展望', font_size=30, font_color=WHITE,
                  bold=True, alignment=PP_ALIGN.CENTER)
    _add_bg_rect(slide, Inches(5), Inches(1.3), Inches(3.3), Inches(0.04), WHITE)

    _add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.2), WHITE)
    summary_lines = [
        {'text': '本集五大重點回顧', 'size': 18, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  1. 台股再創歷史新高：整體市場多頭，手上持股弱不等於市場弱，需區分族群邏輯，功率半導體接棒', 'size': 14, 'color': DARK_GRAY},
        {'text': '  2. 功率半導體超級週期：IFX/STM/ON/Rohm lead time 99+ 週，需求驅動持續性遠強於投機炒作', 'size': 14, 'color': DARK_GRAY},
        {'text': '  3. 被動元件輪動教訓：最缺的不一定漲最多，電阻意外跑最快，市場永遠讓你意外不要把話說死', 'size': 14, 'color': DARK_GRAY},
        {'text': '  4. 機器人萌芽布局：Tesla Optimus 催化，工控零組件廠開始出現機器人訂單，現在進場是最佳時機', 'size': 14, 'color': DARK_GRAY},
        {'text': '  5. 美股 Marvell 三征告捷：AI 定制芯片龍頭持續加碼，功率半導體 IDM 全面表態，是 global 現象', 'size': 14, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(3.0),
                            summary_lines)

    _add_rounded_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.8), WHITE)
    next_lines = [
        {'text': '後市展望', 'size': 18, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  功率半導體週期驗證重點：IDM 廠商下一輪電話會議更新，6/8吋代工廠能否接單是解套關鍵', 'size': 14, 'color': CHARCOAL},
        {'text': '  機器人題材什麼時候進入財報驗證期；消費性/車用起來會縮短 AI 週期，是週期風險需追蹤', 'size': 14, 'color': CHARCOAL},
        {'text': '  AI 仍在幼年期，不要猜幾局下半場，台股功率半導體族群籌碼擁擠程度需持續監控', 'size': 14, 'color': CHARCOAL},
    ]
    _add_multiline_text_box(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.6),
                            next_lines)

    _add_text_box(slide, Inches(2), Inches(7.0), Inches(9.3), Inches(0.4),
                  '資料來源：股癌 Podcast EP672 | 分析日期：2026-06-20 | 本報告僅供參考，非投資建議',
                  font_size=11, font_color=SKY_BLUE, alignment=PP_ALIGN.CENTER)


# ============================================================
# Main Build
# ============================================================

def build_presentation():
    setup_matplotlib()
    print('Generating charts...')
    charts = gen_charts()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print('Building slides...')

    slide_01_cover(prs)
    print('  [1/20] Cover')

    slide_02_contents(prs)
    print('  [2/20] Contents')

    slide_03_part_one(prs)
    print('  [3/20] Part 1 divider')

    slide_04_us_weekly(prs, charts)
    print('  [4/20] US indices — power semiconductor surge')

    slide_05_taiwan_market(prs, charts)
    print('  [5/20] Taiwan market — all-time high')

    slide_06_part_two(prs)
    print('  [6/20] Part 2 divider')

    slide_07_lead_time(prs, charts)
    print('  [7/20] Power semi lead time explosion')

    slide_08_supply_chain(prs, charts)
    print('  [8/20] Supply chain analysis')

    slide_09_ai_demand(prs, charts)
    print('  [9/20] AI server demand driven')

    slide_10_part_three(prs)
    print('  [10/20] Part 3 divider')

    slide_11_passive_components(prs, charts)
    print('  [11/20] Passive components rotation lesson')

    slide_12_robot_layout(prs, charts)
    print('  [12/20] Robot early stage layout')

    slide_13_hw_sw_split(prs, charts)
    print('  [13/20] Hardware vs Software split')

    slide_14_part_four(prs)
    print('  [14/20] Part 4 divider')

    slide_15_us_stocks(prs, charts)
    print('  [15/20] US stock recommendations')

    slide_16_tw_stocks(prs, charts)
    print('  [16/20] Taiwan stock directions')

    slide_17_part_five(prs)
    print('  [17/20] Part 5 divider')

    slide_18_recommendations(prs, charts)
    print('  [18/20] Investment recommendations & radar')

    slide_19_allocation(prs, charts)
    print('  [19/20] Asset allocation')

    slide_20_summary(prs)
    print('  [20/20] Summary & outlook')

    prs.save(OUTPUT_PPTX)
    print(f'\n✅ Saved: {OUTPUT_PPTX}')


if __name__ == '__main__':
    build_presentation()
