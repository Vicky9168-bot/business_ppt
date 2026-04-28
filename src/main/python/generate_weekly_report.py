#!/usr/bin/env python3
"""
早晨財經速解讀 2026-04-28 市場分析報告 — PPT 生成腳本
風格：Corporate Blue & White Business Theme
主題：台積電條款 買盤剛啟動｜外資賣超500億 股市卻創4萬新高
內容：台股衝破4萬點歷史新高、台積電條款散戶搶先、投信需1-3月進場、費半均值回歸資金轉軟體、金管會五箭亞資中心全台落地
分析日期：2026-04-28
"""

import os
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from io import BytesIO
from math import pi

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ============================================================
# Config
# ============================================================

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.abspath(os.path.join(SCRIPT_DIR, '..', '..', '..'))
OUTPUT_DIR = os.path.join(PROJECT_ROOT, 'output', 'reports')
TEMP_DIR = os.path.join(PROJECT_ROOT, 'output', 'temp')
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(TEMP_DIR, exist_ok=True)
OUTPUT_PPTX = os.path.join(OUTPUT_DIR, '20260428_report.pptx')

# Colors — Corporate Blue Theme
CORP_BLUE = RGBColor(0x2B, 0x57, 0x9A)
CORP_BLUE_DARK = RGBColor(0x1A, 0x3A, 0x6C)
CORP_BLUE_LIGHT = RGBColor(0x4A, 0x86, 0xC8)
ACCENT_BLUE = RGBColor(0x5B, 0xA0, 0xD6)
SKY_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
CHARCOAL = RGBColor(0x2D, 0x2D, 0x2D)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x4D)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
BULL_GREEN = RGBColor(0x26, 0xA6, 0x9A)
BEAR_RED = RGBColor(0xEF, 0x53, 0x50)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_CN = 'Microsoft JhengHei'
FONT_EN = 'Arial'
TOTAL_PAGES = 20

# ============================================================
# Matplotlib Setup
# ============================================================

def setup_matplotlib():
    plt.rcParams['font.family'] = ['Heiti TC', 'Arial Unicode MS', 'Hiragino Sans', 'sans-serif']
    plt.rcParams['axes.unicode_minus'] = False
    plt.rcParams['figure.dpi'] = 150
    plt.rcParams['savefig.dpi'] = 250
    plt.rcParams['font.size'] = 14
    plt.rcParams['axes.labelsize'] = 15
    plt.rcParams['axes.titlesize'] = 18


def fig_to_image_stream(fig, dpi=250):
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor='white', edgecolor='none', pad_inches=0.15)
    plt.close(fig)
    buf.seek(0)
    return buf


# ============================================================
# Helper functions
# ============================================================

def _add_bg_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        _set_shape_alpha(shape, alpha)
    return shape


def _set_shape_alpha(shape, alpha_pct):
    fill_elem = shape.fill._fill
    solid = fill_elem.find(qn('a:solidFill'))
    if solid is not None:
        srgb = solid.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_elem = srgb.find(qn('a:alpha'))
            if alpha_elem is None:
                from lxml import etree
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
            alpha_elem.set('val', str(int((100 - alpha_pct) * 1000)))


def _add_text_box(slide, left, top, width, height, text, font_size=14,
                  font_color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name=FONT_CN, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    sp = p._pPr
    if sp is None:
        from lxml import etree
        sp = etree.SubElement(p._p, qn('a:pPr'))
    lnSpc = sp.find(qn('a:lnSpc'))
    if lnSpc is None:
        from lxml import etree
        lnSpc = etree.SubElement(sp, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))
    return txBox


def _add_multiline_text_box(slide, left, top, width, height, lines,
                            font_size=14, font_color=DARK_GRAY,
                            font_name=FONT_CN, line_spacing=1.5,
                            alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_data, dict):
            p.text = line_data.get('text', '')
            p.font.size = Pt(line_data.get('size', font_size))
            p.font.color.rgb = line_data.get('color', font_color)
            p.font.bold = line_data.get('bold', False)
            p.font.name = line_data.get('font', font_name)
        else:
            p.text = str(line_data)
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(2)
    return txBox


def _add_circle(slide, left, top, size, color, text='', font_size=14,
                font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = FONT_EN
        p.alignment = PP_ALIGN.CENTER
    return shape


def _add_rounded_rect(slide, left, top, width, height, color, text='',
                      font_size=13, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(6)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = FONT_CN
        p.alignment = PP_ALIGN.CENTER
    return shape


def _add_header_bar(slide, title_text, subtitle_text=''):
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), CORP_BLUE)
    _add_circle(slide, Inches(0.3), Inches(0.15), Inches(0.8), CORP_BLUE_DARK)
    _add_text_box(slide, Inches(1.3), Inches(0.15), Inches(10), Inches(0.55),
                  title_text, font_size=28, font_color=WHITE, bold=True)
    if subtitle_text:
        _add_text_box(slide, Inches(1.3), Inches(0.65), Inches(10), Inches(0.4),
                      subtitle_text, font_size=15, font_color=SKY_BLUE)
    _add_text_box(slide, Inches(10.5), Inches(0.15), Inches(2.5), Inches(0.4),
                  '早晨財經速解讀', font_size=16, font_color=WHITE, bold=True,
                  alignment=PP_ALIGN.RIGHT)
    _add_text_box(slide, Inches(10.5), Inches(0.55), Inches(2.5), Inches(0.35),
                  '2026.04.28', font_size=12, font_color=SKY_BLUE,
                  alignment=PP_ALIGN.RIGHT)


def _add_footer(slide, page_num, total_pages=TOTAL_PAGES):
    _add_bg_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), LIGHT_GRAY)
    _add_text_box(slide, Inches(0.5), Inches(7.15), Inches(7), Inches(0.3),
                  '早晨財經速解讀 — 2026.04.28 市場分析報告', font_size=10,
                  font_color=MID_GRAY)
    _add_text_box(slide, Inches(10), Inches(7.15), Inches(3), Inches(0.3),
                  f'{page_num} / {total_pages}', font_size=10,
                  font_color=MID_GRAY, alignment=PP_ALIGN.RIGHT, font_name=FONT_EN)


def _add_section_divider(slide, num, title, subtitle, detail=''):
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CORP_BLUE)
    _add_circle(slide, Inches(1.5), Inches(2), Inches(2.5), WHITE, num,
                font_size=48, font_color=CORP_BLUE)
    _add_text_box(slide, Inches(5), Inches(2.2), Inches(7), Inches(1),
                  title, font_size=40, font_color=WHITE, bold=True, font_name=FONT_EN)
    _add_text_box(slide, Inches(5), Inches(3.3), Inches(7), Inches(0.6),
                  subtitle, font_size=18, font_color=SKY_BLUE)
    _add_bg_rect(slide, Inches(5), Inches(4.1), Inches(4), Inches(0.04), WHITE)
    if detail:
        _add_text_box(slide, Inches(5), Inches(4.4), Inches(7), Inches(0.5),
                      detail, font_size=14, font_color=SKY_BLUE)


# ============================================================
# Chart Generation
# ============================================================

def gen_charts():
    charts = {}
    blue_hex = '#2B579A'
    blue_light_hex = '#4A86C8'
    accent_hex = '#5BA0D6'
    green_hex = '#27AE60'
    red_hex = '#E84D4D'
    orange_hex = '#F59E0B'
    purple_hex = '#8E44AD'

    # ── 1. US Indices 4/25 Performance — 標普/那指創新高，費半均值回歸 ──
    fig, ax = plt.subplots(figsize=(11, 3.7))
    indices = ['道瓊工業\n49,620', '標普 500\n7,196', '納斯達克\n24,708', '費城半導體\n10,028']
    daily_chg = [0.52, 0.74, 1.02, -0.48]
    colors = [green_hex if v > 0 else red_hex for v in daily_chg]
    bars = ax.barh(indices, daily_chg, color=colors, height=0.5, edgecolor='white', linewidth=1.5)
    ax.axvline(x=0, color='#ccc', linewidth=1.2, linestyle='-')
    for bar, val in zip(bars, daily_chg):
        x_pos = bar.get_width() + 0.05 if val > 0 else bar.get_width() - 0.05
        ha = 'left' if val > 0 else 'right'
        label = f'+{val:.2f}%' if val > 0 else f'{val:.2f}%'
        ax.text(x_pos, bar.get_y() + bar.get_height() / 2,
                label, va='center', ha=ha, fontsize=16,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xlim(-1.5, 2.0)
    ax.set_title('美股四大指數 — 標普/那指齊創歷史新高，費半均值回歸（2026/04/25）',
                 fontsize=17, fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.tick_params(axis='both', labelsize=14)
    ax.annotate('資金輪向軟體', xy=(-0.48, 3), xytext=(-1.2, 2.5),
                arrowprops=dict(arrowstyle='->', color=red_hex, lw=2),
                fontsize=13, color=red_hex, fontweight='bold')
    fig.tight_layout()
    charts['weekly_us'] = fig_to_image_stream(fig)

    # ── 2. 台股走勢 4/23-4/28 台積電條款拉盤 ──
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4))
    dates = ['4/23', '4/24', '4/25', '4/27', '4/28']
    taiex = [38434, 38734, 38950, 39616, 39610]
    ax1.plot(dates, taiex, 'o-', color=blue_hex, linewidth=2.5, markersize=8)
    ax1.fill_between(range(len(dates)), taiex, min(taiex)*0.99,
                     alpha=0.15, color=blue_hex)
    ax1.set_xticks(range(len(dates)))
    ax1.set_xticklabels(dates, fontsize=11)
    ax1.set_ylabel('點數', fontsize=12, color=blue_hex)
    ax1.set_title('台股指數走勢（4/23-4/28）', fontsize=15, fontweight='bold',
                  color=blue_hex, pad=10)
    ax1.annotate('盤中40,194\n歷史新高！\n收39,616', xy=(3, 39616), xytext=(1.5, 39750),
                 arrowprops=dict(arrowstyle='->', color=green_hex, lw=1.5),
                 fontsize=11, color=green_hex, fontweight='bold')
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.set_ylim(38200, 40500)

    # 外資 vs 投信買賣超
    players = ['外資\n賣超', '投信\n買超', '自營商\n避險', '散戶\n預期買進']
    flows = [-500, 120, -30, 410]
    flow_colors = [red_hex if v < 0 else green_hex for v in flows]
    bars = ax2.bar(range(len(players)), flows, color=flow_colors, width=0.55,
                   edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars, flows):
        y_pos = val + 15 if val >= 0 else val - 30
        label = f'+{val}億' if val > 0 else f'{val}億'
        ax2.text(bar.get_x() + bar.get_width()/2, y_pos, label,
                 ha='center', fontsize=11, fontweight='bold', color=bar.get_facecolor())
    ax2.set_xticks(range(len(players)))
    ax2.set_xticklabels(players, fontsize=11)
    ax2.set_ylabel('買賣超（億元）', fontsize=12, color=blue_hex)
    ax2.set_title('4/27 各法人動向（億元）', fontsize=15, fontweight='bold',
                  color=blue_hex, pad=10)
    ax2.axhline(y=0, color='#ccc', linewidth=1.5)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    fig.tight_layout()
    charts['asia_weekly'] = fig_to_image_stream(fig)

    # ── 3. 資金從硬體轉軟體 (費半整理 vs 七巨頭軟體) ──
    fig, ax = plt.subplots(figsize=(9, 3.8))
    categories = ['費半 SOX\n估值60倍', 'IGV 軟體\nETF', '七巨頭\n（剔除特斯拉）', '標普500\n指數']
    week_chg = [-2.1, 4.5, 3.2, 0.74]
    bar_colors = [green_hex if v > 0 else red_hex for v in week_chg]
    bars = ax.bar(range(len(categories)), week_chg, color=bar_colors,
                  width=0.5, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, week_chg)):
        y_pos = val + 0.2 if val >= 0 else val - 0.5
        label = f'+{val:.1f}%' if val > 0 else f'{val:.1f}%'
        ax.text(i, y_pos, label, ha='center', fontsize=14,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(categories)))
    ax.set_xticklabels(categories, fontsize=12)
    ax.set_ylabel('近期漲跌（%）', fontsize=13, color=blue_hex)
    ax.set_title('資金從硬體轉軟體：費半估值60倍整理，七巨頭低本益比吸資金', fontsize=16,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.axhline(y=0, color='#ccc', linewidth=1)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(-4, 7)
    fig.tight_layout()
    charts['software_pe'] = fig_to_image_stream(fig)

    # ── 4. 資金輪動：AI硬體→AI軟體 ──
    fig, ax = plt.subplots(figsize=(9, 4))
    sectors = ['AI硬體\n(SOX)', 'AI軟體\n(IGV)', '七巨頭\n(MAG6)', '能源\n(XLE)', '消費\n(XLY)']
    fund_flow = [-15, 22, 18, 3, -8]
    flow_colors = [green_hex if v > 0 else red_hex for v in fund_flow]
    bars = ax.bar(range(len(sectors)), fund_flow, color=flow_colors,
                  width=0.5, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, fund_flow)):
        y_pos = val + 1.0 if val >= 0 else val - 2.5
        label = f'+{val}' if val > 0 else f'{val}'
        ax.text(i, y_pos, label, ha='center', fontsize=14,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(sectors)))
    ax.set_xticklabels(sectors, fontsize=12)
    ax.set_ylabel('相對資金流入（相對值）', fontsize=13, color=blue_hex)
    ax.set_title('資金輪動：費半過熱資金外溢，軟體/MAG6 低本益比成接棒方向', fontsize=16,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.axhline(y=0, color='#ccc', linewidth=1)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    fig.tight_layout()
    charts['peg_trend'] = fig_to_image_stream(fig)

    # ── 5. 台積電條款：投信資金流入時程 ──
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4))
    # Left: 預估流入規模
    phases = ['散戶\n（已進場）', '自營商\n（部分）', '投信\n（1-3月後）', '外資\n（待確認）']
    phase_vals = [400, 80, 1700, 0]
    phase_colors = [green_hex, blue_light_hex, orange_hex, '#aaa']
    bars1 = ax1.bar(range(len(phases)), phase_vals, color=phase_colors, width=0.5,
                    edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars1, phase_vals):
        y_pos = val + 30
        label = f'{val}億' if val > 0 else '待觀察'
        ax1.text(bar.get_x() + bar.get_width()/2, y_pos, label,
                 ha='center', fontsize=12, fontweight='bold', color=bar.get_facecolor())
    ax1.set_xticks(range(len(phases)))
    ax1.set_xticklabels(phases, fontsize=12)
    ax1.set_ylabel('預估買入規模（億元）', fontsize=12, color=blue_hex)
    ax1.set_title('台積電條款\n各方資金進場時程', fontsize=14, fontweight='bold',
                  color=blue_hex, pad=10)
    ax1.annotate('合計預估\n1,500-2,000億', xy=(2, 1700), xytext=(0.5, 1600),
                 arrowprops=dict(arrowstyle='->', color=orange_hex, lw=1.5),
                 fontsize=11, color=orange_hex, fontweight='bold')
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.set_ylim(0, 2200)
    # Right: 投信台積電持股比例
    months = ['歷史\n平均', '上限放寬\n前', '目前\n(3.11%)', '上限\n(25%)']
    ratios = [2.2, 2.8, 3.11, 25]
    ratio_colors = [blue_light_hex, blue_light_hex, green_hex, orange_hex]
    bars2 = ax2.bar(range(len(months)), ratios, color=ratio_colors, width=0.5,
                    edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars2, ratios):
        ax2.text(bar.get_x() + bar.get_width()/2, val + 0.3, f'{val}%',
                 ha='center', fontsize=12, fontweight='bold', color=bar.get_facecolor())
    ax2.set_xticks(range(len(months)))
    ax2.set_xticklabels(months, fontsize=12)
    ax2.set_ylabel('台積電持股佔基金比例（%）', fontsize=11, color=blue_hex)
    ax2.set_title('投信台積電持股比例\n歷史新高3.11%，上限達25%', fontsize=13, fontweight='bold',
                  color=blue_hex, pad=10)
    ax2.annotate('空間\n還很大！', xy=(3, 25), xytext=(2.2, 22),
                 arrowprops=dict(arrowstyle='->', color=orange_hex, lw=1.5),
                 fontsize=12, color=orange_hex, fontweight='bold')
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.set_ylim(0, 30)
    fig.tight_layout()
    charts['taiwan_export'] = fig_to_image_stream(fig)

    # ── 6. 美股 IPO 市場：下半年搶錢大作戰 ──
    fig, ax = plt.subplots(figsize=(10, 4.2))
    years = ['2022\n(高利率崩)', '2023\n(稍回溫)', '2024\n(穩步回升)', '2025\n(380億)', '2026年\n(1,600億預)']
    ipo_vals = [25, 120, 210, 380, 1600]
    ipo_colors = [red_hex, blue_light_hex, blue_light_hex, blue_hex, '#1a7a4a']
    bars = ax.bar(range(len(years)), ipo_vals, color=ipo_colors,
                  width=0.55, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, ipo_vals)):
        y_pos = val + 20
        label = f'{val}億' if i < 4 else f'{val}億\n歷史新高!'
        ax.text(i, y_pos, label, ha='center', fontsize=12 if i < 4 else 13,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(years)))
    ax.set_xticklabels(years, fontsize=11)
    ax.set_ylabel('IPO 募資金額（億美元）', fontsize=13, color=blue_hex)
    ax.set_title('美股IPO市場全面回歸：2026年預估1,600億！下半年SpaceX/Anthropic/OpenAI搶錢', fontsize=16,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    fig.tight_layout()
    charts['capex_war'] = fig_to_image_stream(fig)

    # ── 7. 產業動能雷達圖 ──
    fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
    categories_r = ['AI 硬體\n(費半)', '台積電\n條款', '金管會\n五箭', '軟體股\n回升', 'IPO\n搶錢']
    N = len(categories_r)
    current = [7, 8, 9, 6, 7]
    future = [7, 9, 9, 8, 9]
    angles = [n / float(N) * 2 * pi for n in range(N)]
    angles += angles[:1]
    current += current[:1]
    future += future[:1]
    ax.plot(angles, current, 'o-', color=blue_hex, linewidth=2.5,
            label='當前評分', markersize=8)
    ax.fill(angles, current, alpha=0.15, color=blue_hex)
    ax.plot(angles, future, 's--', color=orange_hex, linewidth=2.5,
            label='3-6月展望', markersize=8)
    ax.fill(angles, future, alpha=0.1, color=orange_hex)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories_r, fontsize=12, fontweight='bold')
    ax.set_ylim(0, 10)
    ax.set_yticks([2, 4, 6, 8, 10])
    ax.set_yticklabels(['2', '4', '6', '8', '10'], fontsize=10, color='gray')
    ax.legend(fontsize=12, loc='upper right', bbox_to_anchor=(1.3, 1.1))
    ax.set_title('市場動能雷達圖', fontsize=18,
                 fontweight='bold', color=blue_hex, pad=30)
    fig.tight_layout()
    charts['sector_radar'] = fig_to_image_stream(fig)

    # ── 8. 金管會五箭影響力 ──
    fig, ax = plt.subplots(figsize=(9, 3.5))
    policies = ['亞資中心\n全台落地(7月)', '台積電\n持股25%', '虛擬資產\n管理條例', '0股交易\n整股同步', '家族辦公室\n私募開放']
    impact = [9, 9, 7, 6, 8]
    p_colors = ['#1a7a4a', blue_hex, blue_light_hex, orange_hex, accent_hex]
    bars = ax.bar(range(len(policies)), impact, color=p_colors,
                  width=0.55, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, impact)):
        ax.text(i, val + 0.2, f'{val}/10', ha='center', fontsize=13,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(policies)))
    ax.set_xticklabels(policies, fontsize=10)
    ax.set_ylabel('政策影響力（1-10）', fontsize=13, color=blue_hex)
    ax.set_title('金管會五箭齊發：亞資中心7月全台落地，制度×場域×產品三線並進', fontsize=16,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(0, 11.5)
    fig.tight_layout()
    charts['capex_ratio'] = fig_to_image_stream(fig)

    # ── 9. 資產配置建議 ──
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(8.5, 3.5))
    labels_c = ['台積電/AI硬體', '主動型ETF', '現金', '軟體股/MAG6', '其他']
    sizes_c = [30, 25, 25, 10, 10]
    colors_c = [blue_hex, blue_light_hex, '#95a5a6', orange_hex, accent_hex]
    wedges, texts, autotexts = ax1.pie(
        sizes_c, labels=labels_c, autopct='%1.0f%%', colors=colors_c,
        startangle=90, pctdistance=0.75, textprops={'fontsize': 11})
    for at in autotexts:
        at.set_fontweight('bold')
        at.set_color('white')
    ax1.set_title('積極型（等待投信進場）', fontsize=14, fontweight='bold', color=blue_hex, pad=15)

    labels_g = ['現金', '台積電', '主動ETF', '債券/黃金', '其他']
    sizes_g = [40, 25, 15, 10, 10]
    colors_g = ['#95a5a6', blue_hex, blue_light_hex, orange_hex, accent_hex]
    wedges2, texts2, autotexts2 = ax2.pie(
        sizes_g, labels=labels_g, autopct='%1.0f%%', colors=colors_g,
        startangle=90, pctdistance=0.75, textprops={'fontsize': 11})
    for at in autotexts2:
        at.set_fontweight('bold')
        at.set_color('white')
    ax2.set_title('穩健型（等台積電條款落地）', fontsize=14, fontweight='bold', color=blue_hex, pad=15)
    fig.tight_layout()
    charts['allocation'] = fig_to_image_stream(fig)

    # ── 10. 外資賣超 vs 台股走勢（7年對比）──
    fig, ax = plt.subplots(figsize=(11, 4.0))
    years = ['2020', '2021', '2022', '2023', '2024', '2025', '2026\nYTD']
    foreign_sell = [-5300, -4500, -12000, 0, -6900, -5900, -4800]
    bar_colors = [red_hex if v < 0 else green_hex for v in foreign_sell]
    bars = ax.bar(range(len(years)), foreign_sell, color=bar_colors,
                  width=0.55, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, foreign_sell)):
        y_pos = val - 350 if val < 0 else val + 100
        label = f'{val:,}億' if val != 0 else '小買'
        ax.text(i, y_pos, label, ha='center', fontsize=11,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(years)))
    ax.set_xticklabels(years, fontsize=12)
    ax.set_ylabel('外資買賣超（億元）', fontsize=13, color=blue_hex)
    ax.set_title('外資7年賣超台股：台股從8,000→40,000點翻5倍，外資回補是最大潛在觸媒', fontsize=16,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.axhline(y=0, color='#ccc', linewidth=1.5)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    fig.tight_layout()
    charts['asset_ytd'] = fig_to_image_stream(fig)

    # ── 11. 主動型ETF台積電持股比例（前三名）──
    fig, ax = plt.subplots(figsize=(10, 3.5))
    etf_names = ['野村台灣\n00985A', '富華未來\n500B', '第一金台股優\n00994A', '992A', '981A', '0050\n台灣50']
    tsmc_weight = [29.0, 14.93, 4.575, 6.2, 7.8, 10.5]
    etf_colors = ['#1a7a4a', '#1a7a4a', green_hex, blue_hex, blue_light_hex, '#aaa']
    bars = ax.bar(range(len(etf_names)), tsmc_weight, color=etf_colors,
                  width=0.55, edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars, tsmc_weight)):
        y_pos = val + 0.5
        ax.text(i, y_pos, f'{val}%', ha='center', fontsize=12,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_xticks(range(len(etf_names)))
    ax.set_xticklabels(etf_names, fontsize=11)
    ax.set_ylabel('台積電持股比例（%）', fontsize=13, color=blue_hex)
    ax.set_title('持有台積電最高的主動ETF：野村00985A以29%奪冠，條款鬆綁後更大受益', fontsize=16,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_ylim(0, 36)
    ax.axhline(y=10, color=orange_hex, linewidth=1.5, linestyle='--', alpha=0.7)
    ax.text(5.4, 10.3, '舊上限10%', fontsize=10, color=orange_hex)
    ax.axhline(y=25, color=green_hex, linewidth=1.5, linestyle='--', alpha=0.7)
    ax.text(5.4, 25.3, '新上限25%', fontsize=10, color=green_hex)
    fig.tight_layout()
    charts['global_stocks_ytd'] = fig_to_image_stream(fig)

    # ── 12. 台灣景氣 & 富豪客戶暴增 ──
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4.0))
    # 景氣燈號趨勢
    months_jq = ['2025\n10月', '2025\n12月', '2026\n1月', '2026\n2月', '2026\n3月']
    scores = [42, 41, 40, 41, 39]
    light_colors = ['#E84D4D']*5
    ax1.bar(range(len(months_jq)), scores, color=light_colors, width=0.5,
            edgecolor='white', linewidth=1.5)
    for i, val in enumerate(scores):
        ax1.text(i, val + 0.3, f'{val}分', ha='center', fontsize=12,
                 fontweight='bold', color='#E84D4D')
    ax1.axhline(y=38, color='#aaa', linewidth=2, linestyle='--', alpha=0.7)
    ax1.text(4.5, 38.2, '紅燈下限38', fontsize=10, color='gray')
    ax1.set_xticks(range(len(months_jq)))
    ax1.set_xticklabels(months_jq, fontsize=11)
    ax1.set_ylabel('景氣燈號分數', fontsize=12, color=blue_hex)
    ax1.set_title('台灣景氣燈號\n紅燈持續（3月39分）', fontsize=14, fontweight='bold',
                  color=blue_hex, pad=10)
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.set_ylim(35, 46)
    # 富豪客戶 AUM
    aum_years = ['2024\nAUM', '2025\nAUM', '2026預\nAUM']
    aum_vals = [1.1, 1.55, 3.0]
    aum_colors = [blue_light_hex, blue_hex, '#1a7a4a']
    bars2 = ax2.bar(range(len(aum_years)), aum_vals, color=aum_colors, width=0.5,
                    edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars2, aum_vals):
        ax2.text(bar.get_x() + bar.get_width()/2, val + 0.05, f'{val}兆',
                 ha='center', fontsize=13, fontweight='bold', color=bar.get_facecolor())
    ax2.set_xticks(range(len(aum_years)))
    ax2.set_xticklabels(aum_years, fontsize=12)
    ax2.set_ylabel('資產管理規模（兆元）', fontsize=12, color=blue_hex)
    ax2.set_title('台灣高資產客戶AUM\n年增55%，今年預估突破3兆', fontsize=14, fontweight='bold',
                  color=blue_hex, pad=10)
    ax2.annotate('突破3兆！\n客戶人數+72%', xy=(2, 3.0), xytext=(1.2, 2.7),
                 arrowprops=dict(arrowstyle='->', color='#1a7a4a', lw=1.5),
                 fontsize=11, color='#1a7a4a', fontweight='bold')
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.set_ylim(0, 3.8)
    fig.tight_layout()
    charts['commodity_ytd'] = fig_to_image_stream(fig)

    # ── 13. 台積電條款效應：三股力量共振 ──
    fig, ax = plt.subplots(figsize=(10, 3.8))
    forces = ['台積電EPS\n2年翻倍', '投信持股\n上限鬆綁', 'ETF資金\n定期扣款', '外資\n潛在回補', '政策\n亞資中心']
    force_score = [9, 9, 8, 7, 8]
    f_colors = [blue_hex, '#1a7a4a', blue_light_hex, orange_hex, accent_hex]
    bars = ax.barh(range(len(forces)), force_score, color=f_colors,
                   height=0.55, edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars, force_score):
        x_pos = bar.get_width() + 0.1
        ax.text(x_pos, bar.get_y() + bar.get_height() / 2,
                f'{val}/10', va='center', ha='left', fontsize=12,
                fontweight='bold', color=bar.get_facecolor())
    ax.set_yticks(range(len(forces)))
    ax.set_yticklabels(forces, fontsize=12)
    ax.set_xlabel('影響力評分（1-10）', fontsize=13, color=blue_hex)
    ax.set_title('支撐台股4萬點的五大力量：基本面×政策×資金三線共振', fontsize=16,
                 fontweight='bold', color=blue_hex, pad=15)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color('#ddd')
    ax.spines['left'].set_color('#ddd')
    ax.set_xlim(0, 12)
    ax.invert_yaxis()
    fig.tight_layout()
    charts['sector_rotation'] = fig_to_image_stream(fig)

    # ── 14. 台股結構解析：一個人的50 ──
    fig, ax = plt.subplots(figsize=(10, 3.5))
    stock_cat = ['台積電\n（拉指數）', '0050等\n大型ETF', '中大型\n權值股', '中小型股\n（下跌）', '總體\n下跌檔數']
    stock_change = [3.2, 1.8, -1.2, -4.5, 1300]
    # Use two axes
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 3.8))
    cat1 = ['台積電\n（拉指數）', '0050等\n大型ETF', '中大型\n權值股', '中小型股\n（下跌）']
    chg1 = [3.2, 1.8, -1.2, -4.5]
    s_colors = [green_hex, blue_light_hex, red_hex, '#8B0000']
    bars1 = ax1.bar(range(len(cat1)), chg1, color=s_colors, width=0.5,
                    edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars1, chg1)):
        y_pos = val + 0.2 if val >= 0 else val - 0.5
        label = f'+{val:.1f}%' if val > 0 else f'{val:.1f}%'
        ax1.text(i, y_pos, label, ha='center', fontsize=13,
                 fontweight='bold', color=bar.get_facecolor())
    ax1.set_xticks(range(len(cat1)))
    ax1.set_xticklabels(cat1, fontsize=11)
    ax1.set_ylabel('單日漲跌幅（%）', fontsize=12, color=blue_hex)
    ax1.set_title('一個人的50：台積電漲\n其餘1300檔下跌', fontsize=14, fontweight='bold',
                  color=blue_hex, pad=10)
    ax1.axhline(y=0, color='#ccc', linewidth=1)
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.set_ylim(-6, 5)
    # Right: 主動ETF YTD
    etf2 = ['994A', '981A', '992A', '主動ETF\n均值', '大盤\n加權指數']
    ytd2 = [65, 64, 62, 44, 25]
    etf2_colors = ['#1a7a4a', '#1a7a4a', '#1a7a4a', green_hex, '#aaa']
    bars2 = ax2.bar(range(len(etf2)), ytd2, color=etf2_colors, width=0.5,
                    edgecolor='white', linewidth=1.5)
    for i, (bar, val) in enumerate(zip(bars2, ytd2)):
        ax2.text(i, val + 1.5, f'+{val}%', ha='center', fontsize=12,
                 fontweight='bold', color=bar.get_facecolor())
    ax2.set_xticks(range(len(etf2)))
    ax2.set_xticklabels(etf2, fontsize=11)
    ax2.set_ylabel('YTD 漲幅（%）', fontsize=12, color=blue_hex)
    ax2.set_title('主動ETF全面跑贏大盤\n3檔突破60%', fontsize=14, fontweight='bold',
                  color=blue_hex, pad=10)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.set_ylim(0, 78)
    fig.tight_layout()
    charts['sp500_earnings'] = fig_to_image_stream(fig)

    return charts


# ============================================================
# Slide Builders
# ============================================================

def slide_01_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CORP_BLUE)
    _add_circle(slide, Inches(10.5), Inches(-0.5), Inches(3), CORP_BLUE_DARK)
    _add_circle(slide, Inches(11.5), Inches(5), Inches(2.5), CORP_BLUE_DARK)
    _add_circle(slide, Inches(-0.8), Inches(5.5), Inches(2), CORP_BLUE_DARK)
    _add_rounded_rect(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(5), WHITE)
    _add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
                  '早晨財經速解讀 — 市場分析報告', font_size=38, font_color=CORP_BLUE,
                  bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(0.9),
                  '台積電條款 買盤剛啟動｜外資賣超500億 股市卻創4萬新高',
                  font_size=19, font_color=CORP_BLUE_DARK,
                  bold=True, alignment=PP_ALIGN.CENTER)
    _add_bg_rect(slide, Inches(4.5), Inches(3.5), Inches(4.3), Inches(0.05), CORP_BLUE)
    _add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.5),
                  '台股盤中40,194點歷史新高 | 投信需1-3月進場 | 費半估值60倍均值回歸',
                  font_size=14, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10), Inches(0.5),
                  '標普500&那指齊創新高 | 金管會五箭亞資中心7月全台落地 | IPO下半年搶錢1,600億',
                  font_size=14, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.4),
                  '2026 年 4 月 28 日', font_size=12,
                  font_color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.5),
                  '早晨財經速解讀 — 2026.04.28 綜合分析', font_size=18,
                  font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


def slide_02_contents(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '目錄 CONTENTS', '2026.04.28 台積電條款+台股4萬+費半整理+金管會五箭+IPO搶錢')
    _add_footer(slide, 2)

    sections = [
        {'num': '01', 'title': '市場總覽', 'desc': '台股衝破4萬點\n外資賣500億仍創高\n一個人的50現象', 'color': CORP_BLUE},
        {'num': '02', 'title': '台積電條款', 'desc': '散戶先搶、投信待進\n1,500-2,000億流入\n3個月進場時程', 'color': CORP_BLUE_LIGHT},
        {'num': '03', 'title': '台灣金融', 'desc': '景氣紅燈39分\n富豪客戶+72%\n金管會五箭全台落地', 'color': ACCENT_BLUE},
        {'num': '04', 'title': '投資策略', 'desc': '擁擠交易部署\nIPO搶錢應對\n後市觀察重點', 'color': ACCENT_PURPLE},
    ]

    start_x = Inches(0.8)
    for i, sec in enumerate(sections):
        x = start_x + i * Inches(3.2)
        y = Inches(2.0)
        _add_rounded_rect(slide, x, y, Inches(2.9), Inches(4.3), WHITE)
        _add_circle(slide, x + Inches(0.85), y + Inches(0.3), Inches(1.1),
                    sec['color'], sec['num'], font_size=28)
        _add_text_box(slide, x + Inches(0.1), y + Inches(1.6), Inches(2.7),
                      Inches(0.5), sec['title'], font_size=20,
                      font_color=CHARCOAL, bold=True, alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, x + Inches(0.2), y + Inches(2.2), Inches(2.5),
                      Inches(1.5), sec['desc'], font_size=13,
                      font_color=MID_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.5)
        _add_bg_rect(slide, x + Inches(0.8), y + Inches(4.0), Inches(1.3),
                     Inches(0.05), sec['color'])


def slide_03_part_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '01', 'PART ONE.', '市場總覽 — 台股衝破4萬點',
                         '盤中40,194歷史新高 | 外資賣500億 內資撐盤 | 標普&那指創新高')


def slide_04_us_weekly(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '美股四大指數 — 標普/那指創新高，費半均值回歸', '2026/04/25 — 標普+0.74% 那指+1.02% 費半-0.48%（估值60倍整理）')
    _add_footer(slide, 4)

    idx_data = [
        ('道瓊工業', '49,620', '+252 pt', '+0.52%', ACCENT_GREEN),
        ('標普 500', '7,196', '+53 pt', '+0.74%', ACCENT_GREEN),
        ('納斯達克', '24,708', '+249 pt', '+1.02%', ACCENT_GREEN),
        ('費城半導體', '10,028', '-49 pt', '-0.48%', ACCENT_RED),
    ]
    for i, (name, price, change, pct, color) in enumerate(idx_data):
        x = Inches(0.4) + i * Inches(3.2)
        _add_rounded_rect(slide, x, Inches(1.4), Inches(3.0), Inches(1.6), WHITE)
        _add_bg_rect(slide, x + Inches(0.05), Inches(1.4), Inches(2.9), Inches(0.06), color)
        _add_text_box(slide, x + Inches(0.2), Inches(1.55), Inches(2.6), Inches(0.35),
                      name, font_size=15, font_color=CHARCOAL, bold=True, alignment=PP_ALIGN.CENTER)
        _add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(2.6), Inches(0.4),
                      price, font_size=24, font_color=CHARCOAL, bold=True,
                      alignment=PP_ALIGN.CENTER, font_name=FONT_EN)
        _add_text_box(slide, x + Inches(0.2), Inches(2.35), Inches(2.6), Inches(0.35),
                      f'{change}（{pct}）', font_size=14, font_color=color,
                      bold=True, alignment=PP_ALIGN.CENTER)

    img = charts['weekly_us']
    slide.shapes.add_picture(img, Inches(1.2), Inches(3.2), Inches(11), Inches(3.7))


def slide_05_taiwan_market(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '台股衝破4萬點 — 歷史新高，但結構不健康', '4/27 盤中40,194歷史新高 | 收39,616(+684pt) | 外資賣500億 | 1,300檔下跌')
    _add_footer(slide, 5)

    img = charts['asia_weekly']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.8), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.3), Inches(1.3), Inches(4.7), Inches(3.5), WHITE)
    taiwan_lines = [
        {'text': '4萬點的真相', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  盤中40,194點 → 歷史新高', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  收盤39,616（+684pt）日K收黑', 'size': 14, 'color': ACCENT_RED, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  外資賣超500億（數字大）', 'size': 13, 'color': ACCENT_RED},
        {'text': '  投信買超百億（主力撐盤）', 'size': 13, 'color': ACCENT_GREEN},
        {'text': '  尾盤600億元賣壓壓回', 'size': 13, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  1,300檔下跌 = 一個人的50', 'size': 13, 'color': ACCENT_RED},
        {'text': '  今年第31次創歷史新高', 'size': 12, 'color': MID_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(8.5), Inches(1.45), Inches(4.3), Inches(3.2),
                            taiwan_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    bottom = [
        {'text': '三股力量共振的結果', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  ① 台積電條款預期效應：散戶搶先投信進場 → 推升台積電 → 指數被拉上4萬', 'size': 14, 'color': DARK_GRAY},
        {'text': '  ② ETF定期扣款：每月固定資金流入撐盤 ③ 政策鬆綁：金管會持續開大門', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), bottom)


def slide_06_part_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '02', 'PART TWO.', '台積電條款 — 散戶先行，投信待動',
                         '1,500-2,000億流入預估 | 投信需1-3個月 | 散戶是先行部隊')


def slide_07_tsmc_clause(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '台積電條款：買盤剛啟動，大部隊還沒進場', '投信持股3.11%歷史新高 | 契約修改→董事會→主管機關：至少1-3個月')
    _add_footer(slide, 7)

    img = charts['taiwan_export']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.8))

    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(4.9), Inches(5.5), WHITE)
    tsmc_lines = [
        {'text': '台積電條款三階段', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  第一階段（現在）：散戶搶先', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  「我知道你要買，我先買給你」', 'size': 13, 'color': DARK_GRAY},
        {'text': '  預期交易推升急漲行情', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  第二階段（1-3個月後）：投信進場', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  需修改契約→董事會通過', 'size': 13, 'color': DARK_GRAY},
        {'text': '  →報主管機關審核', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  第三階段（待確認）：外資觀察', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  新台幣同步走強是外資訊號', 'size': 13, 'color': DARK_GRAY},
        {'text': '  一度突破31.4元（兩月新高）', 'size': 13, 'color': ACCENT_GREEN},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.5), Inches(5.3),
                            tsmc_lines)


def slide_08_active_etf_tsmc(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '主動ETF台積電持股排行 — 條款鬆綁後誰受益最大', '野村00985A持股29% | 舊上限10% | 新上限25% | ETF永動機現象消失')
    _add_footer(slide, 8)

    img = charts['global_stocks_ytd']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(8), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.6), Inches(1.3), Inches(4.4), Inches(3.5), WHITE)
    etf_lines = [
        {'text': '持股上限鬆綁效應', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  野村00985A：29%（已超新上限）', 'size': 13, 'color': ACCENT_RED, 'bold': True},
        {'text': '  → 已超過25%，須關注調整', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  持股10-25%的ETF最受益', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  可直接增持至25%上限', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  ETF永動機消失', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  以前台積超10%→被迫買0050', 'size': 12, 'color': DARK_GRAY},
        {'text': '  現在可直接配置台積電', 'size': 12, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(8.8), Inches(1.45), Inches(4.0), Inches(3.2),
                            etf_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    bottom = [
        {'text': '今年主動ETF表現（YTD）', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  994A +65% | 981A +64% | 992A +62% ← 持股集中AI/台積電，貝塔值放大效應', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  主動ETF均值+44% vs 大盤+25%。風水輪流轉：高股息→全市值→主動型，每年都有新主流', 'size': 13, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), bottom)


def slide_09_fund_rotation(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '資金輪動 — 費半均值回歸，資金轉向軟體/MAG6', '費半估值60倍 | 七巨頭（剔除特斯拉）低於5年均值 | 下半年財報季觸媒')
    _add_footer(slide, 9)

    img = charts['software_pe']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.5))

    img2 = charts['peg_trend']
    slide.shapes.add_picture(img2, Inches(0.3), Inches(5.0), Inches(7.5), Inches(2.3))

    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(5.0), Inches(6.0), WHITE)
    rot_lines = [
        {'text': '輪動邏輯解析', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  費半現況：估值60倍', 'size': 14, 'color': ACCENT_RED, 'bold': True},
        {'text': '  連18天上漲 = 籌碼擁擠', 'size': 13, 'color': DARK_GRAY},
        {'text': '  最佳出路：盤整讓時間消化', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  七巨頭機會：低本益比', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  Alphabet/微軟/Amazon/Meta', 'size': 13, 'color': DARK_GRAY},
        {'text': '  低於5年平均本益比', 'size': 13, 'color': DARK_GRAY},
        {'text': '  陸續公布業績 = 觸媒', 'size': 13, 'color': ACCENT_GREEN},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  整體標普500：均值24倍', 'size': 13, 'color': CORP_BLUE, 'bold': True},
        {'text': '  合理不貴，資金仍充足', 'size': 12, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  YTD標普500 +5.1%', 'size': 13, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  伊朗戰事跌幅僅7.1%', 'size': 12, 'color': DARK_GRAY},
        {'text': '  遠小於2025關稅戰15%+', 'size': 12, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.6), Inches(5.8),
                            rot_lines)


def slide_10_part_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '03', 'PART THREE.', '台灣金融 — 全面升溫',
                         '景氣紅燈39分 | 富豪客戶+72% | 金管會五箭亞資中心全台落地')


def slide_11_taiwan_economy(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '台灣景氣 & 財富效應 — 高檔紅燈 + 富豪暴增', '3月景氣39分紅燈 | 高資產客戶+72% | AUM 2.4兆 | 財富效應擴大消費')
    _add_footer(slide, 11)

    img = charts['commodity_ytd']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.8))

    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(4.9), Inches(5.5), WHITE)
    econ_lines = [
        {'text': '台灣財富效應全面擴散', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  景氣：3月39分，高檔紅燈', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  高資產客戶：23,000人', 'size': 13, 'color': DARK_GRAY},
        {'text': '  年增幅72%！「科技新貴暴發」', 'size': 13, 'color': ACCENT_GREEN},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  AUM 2.4兆，年增55%', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  今年預估突破3兆', 'size': 13, 'color': DARK_GRAY},
        {'text': '  存款比例仍42%（解定存空間大）', 'size': 13, 'color': ACCENT_ORANGE},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  MSCI台灣PB比5.2倍', 'size': 13, 'color': ACCENT_RED, 'bold': True},
        {'text': '  極端值+4標準差', 'size': 12, 'color': DARK_GRAY},
        {'text': '  但ROE也突破20%（2標準差）', 'size': 12, 'color': ACCENT_GREEN},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.5), Inches(5.3),
                            econ_lines)


def slide_12_fsc_policy(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '金管會五箭齊發 — 亞資中心7月全台落地', '全齡金融×信任金融×亞資中心×國際資本×三軌金融 | 從踩剎車→踩油門')
    _add_footer(slide, 12)

    img = charts['capex_ratio']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(5.0), Inches(3.5), WHITE)
    fsc_lines = [
        {'text': '金管會五箭重點', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  ① 亞資中心全台落地（7月）', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  高雄試辦一年→全國推廣', 'size': 13, 'color': DARK_GRAY},
        {'text': '  高雄20家銀行，AUM 3,700億', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  ② 虛擬資產管理條例（下半年）', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  台版穩定幣架構啟動', 'size': 13, 'color': DARK_GRAY},
        {'text': '  泰達/Circle等國際業者可拿牌', 'size': 13, 'color': DARK_GRAY},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  ③ 台灣全球自由度第5名', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  79.8分，前面是新/瑞/愛/澳', 'size': 13, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.6), Inches(3.2),
                            fsc_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    insight = [
        {'text': 'KPI = 鬆綁了多少法規：金融市場發展創新處', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  金管會設立專屬部門，彭主委每月親自盯「鬆綁進度」，從查違規→找機會，台灣金融業從特許產業→國家產業引擎', 'size': 14, 'color': DARK_GRAY},
        {'text': '  0股交易2027年整股同步開盤 | 家族辦公室/境外私募 優先放行 | TISA長期投資帳戶建立', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), insight)


def slide_13_ipo_market(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, 'IPO市場：2026年1,600億歷史新高，下半年搶錢大作戰', 'SpaceX/Anthropic/OpenAI | 目前僅募140-150億 | 資金胃納量考驗')
    _add_footer(slide, 13)

    img = charts['capex_war']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.5))

    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(5.0), Inches(3.5), WHITE)
    ipo_lines = [
        {'text': 'IPO搶錢大作戰', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  2026全年預估：1,600億', 'size': 14, 'color': ACCENT_GREEN, 'bold': True},
        {'text': '  YTD已募：僅140-150億', 'size': 13, 'color': ACCENT_RED},
        {'text': '  → 後面還有1,450億要搶！', 'size': 13, 'color': ACCENT_RED, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  主角：SpaceX（最受期待）', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  Anthropic、OpenAI 待定', 'size': 13, 'color': DARK_GRAY},
        {'text': '  → 市場從防禦→成長敘事', 'size': 13, 'color': ACCENT_GREEN},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  核心問題：資金胃納量？', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
        {'text': '  股價漲 + IPO搶錢 能並存？', 'size': 12, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.6), Inches(3.2),
                            ipo_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    note = [
        {'text': '美台市場解定存浪潮：資金量充沛', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  不只台股游資多，美股也是大量解定存湧入。解定存→股市＋IPO，資金鏈條完整。關鍵是下半年資金夠不夠一邊撐行情一邊被IPO吸走', 'size': 14, 'color': DARK_GRAY},
        {'text': '  → 台股遊資估算：若存款42%的高資產客戶每降1%，流入量達240億，對台股影響不可忽視', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), note)


def slide_14_part_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '04', 'PART FOUR.', '投資策略 — 擁擠交易中的部署',
                         '台積電條款分批進場 | 外資觀察 | 財富效應評估 | 風險矩陣')


def slide_15_market_forces(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '台股4萬點的五大支撐力量', '台積電EPS 2年翻倍 | 投信持股上限 | ETF扣款 | 外資潛力 | 金管會政策')
    _add_footer(slide, 15)

    img = charts['sector_rotation']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(7.5), Inches(3.8))

    _add_rounded_rect(slide, Inches(8.1), Inches(1.3), Inches(4.9), Inches(5.5), WHITE)
    force_lines = [
        {'text': '三根柱子：難以鬆動', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  ① 權值股EPS', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  台積電EPS 1,300→2,200（2年翻倍）', 'size': 13, 'color': DARK_GRAY},
        {'text': '  高盛被迫上調目標價', 'size': 13, 'color': ACCENT_GREEN},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  ② ETF資金', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  下個月突然不扣款？不可能', 'size': 13, 'color': DARK_GRAY},
        {'text': '  每月固定流入是壓艙石', 'size': 13, 'color': ACCENT_GREEN},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  ③ 政策鬆綁', 'size': 14, 'color': CORP_BLUE, 'bold': True},
        {'text': '  突然緊縮？幾乎不可能', 'size': 13, 'color': DARK_GRAY},
        {'text': '  亞資中心才剛要全台落地', 'size': 13, 'color': ACCENT_GREEN},
    ]
    _add_multiline_text_box(slide, Inches(8.3), Inches(1.45), Inches(4.5), Inches(5.3),
                            force_lines)


def slide_16_foreign_vs_domestic(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '扶積滅洋 — 外資7年賣台股，台股翻5倍', '2020-2026 外資有6年賣超 | 台股8,000→40,000 | 外資回補是最大上行觸媒')
    _add_footer(slide, 16)

    img = charts['asset_ytd']
    slide.shapes.add_picture(img, Inches(0.2), Inches(1.3), Inches(6.8), Inches(3.5))

    img2 = charts['sp500_earnings']
    slide.shapes.add_picture(img2, Inches(7.2), Inches(1.3), Inches(5.8), Inches(3.5))

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), WHITE)
    note_lines = [
        {'text': '內資主導的結構性牛市', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  外資賣超≠台股不漲：7年6賣，台股翻5倍。內資ETF扣款+主動型ETF+台積電條款 = 強大內生動力', 'size': 14, 'color': DARK_GRAY},
        {'text': '  外資觀察指標：新台幣匯率突破31.4元（兩月新高） = 外資有流入機會，股匯雙漲是最強訊號', 'size': 14, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3),
                            note_lines)


def slide_17_part_five(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_section_divider(slide, '05', 'PART FIVE.', '投資策略總結 — 2026.04.28 操作建議',
                         '台積電條款分批布局 | 費半整理期策略 | 後市觀察重點')


def slide_18_recommendations(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '產業動能雷達 & 配置建議', 'AI硬體整理等待 | 台積電條款分批進場 | 外資觀察是關鍵')
    _add_footer(slide, 18)

    img = charts['sector_radar']
    slide.shapes.add_picture(img, Inches(0.2), Inches(1.3), Inches(5.2), Inches(5.2))

    _add_rounded_rect(slide, Inches(5.6), Inches(1.3), Inches(7.4), Inches(2.3), WHITE)
    tw_rec_lines = [
        {'text': '積極型配置建議', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  台積電/相關ETF 30%（條款效應，分批布局）', 'size': 13, 'color': CHARCOAL},
        {'text': '  主動型ETF 25%（994A/981A/992A 持續跑贏）', 'size': 13, 'color': CHARCOAL},
        {'text': '  現金 25%（等待投信真正進場確認信號）', 'size': 13, 'color': CHARCOAL},
        {'text': '  軟體股/MAG6 10%（低本益比財報季觸媒）', 'size': 13, 'color': CHARCOAL},
    ]
    _add_multiline_text_box(slide, Inches(5.8), Inches(1.4), Inches(7), Inches(2.1),
                            tw_rec_lines)

    _add_rounded_rect(slide, Inches(5.6), Inches(3.8), Inches(7.4), Inches(3.0), WHITE)
    us_rec_lines = [
        {'text': '穩健型配置建議', 'size': 17, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  現金 40%（高PB 5.2倍，等確認訊號）', 'size': 13, 'color': CHARCOAL},
        {'text': '  台積電/大型ETF 25%（0050/主動型）', 'size': 13, 'color': CHARCOAL},
        {'text': '  防禦型資產 15%（債券/黃金）', 'size': 13, 'color': CHARCOAL},
        {'text': '  其他 10%', 'size': 13, 'color': CHARCOAL},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  注意：台股日K收黑 + 尾盤600億賣壓', 'size': 11, 'color': ACCENT_ORANGE, 'bold': True},
    ]
    _add_multiline_text_box(slide, Inches(5.8), Inches(3.9), Inches(7), Inches(2.8),
                            us_rec_lines)


def slide_19_allocation(prs, charts):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, LIGHT_GRAY)
    _add_header_bar(slide, '資產配置圖解 & 後市觀察重點', '2026.04.28 — 台積電條款買盤剛啟動，分批進場勝於追高')
    _add_footer(slide, 19)

    img = charts['allocation']
    slide.shapes.add_picture(img, Inches(0.3), Inches(1.3), Inches(8.5), Inches(3.5))

    _add_rounded_rect(slide, Inches(9.0), Inches(1.3), Inches(4.0), Inches(3.5), WHITE)
    risk_lines = [
        {'text': '後市觀察重點', 'size': 16, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  1. 投信何時真正進場？', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  契約修改→董事會→主管機關', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  2. 新台幣是否持續走強？', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  匯率強 = 外資流入訊號', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  3. 七巨頭財報週結果', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  軟體資金能否接棒硬體', 'size': 12, 'color': MID_GRAY},
        {'text': '', 'size': 4, 'color': DARK_GRAY},
        {'text': '  4. ETF贖回潮是否出現', 'size': 13, 'color': CHARCOAL, 'bold': True},
        {'text': '  贖回潮 = 本輪系統性壓力', 'size': 12, 'color': MID_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(9.2), Inches(1.4), Inches(3.6), Inches(3.3),
                            risk_lines)

    _add_rounded_rect(slide, Inches(0.3), Inches(5.1), Inches(12.7), Inches(1.7), CORP_BLUE)
    summary = [
        {'text': '核心策略總結', 'size': 18, 'color': WHITE, 'bold': True},
        {'text': '', 'size': 5, 'color': WHITE},
        {'text': '  台積電條款：散戶先行，投信1-3個月後才真正進場。分批布局優於追高。ETF贖回潮未出現前，多頭格局持續', 'size': 15, 'color': WHITE},
        {'text': '  外資賣500億台股仍創高 = 內資力量強大。新台幣走強是外資回補前兆，是最大上行觸媒', 'size': 14, 'color': SKY_BLUE},
    ]
    _add_multiline_text_box(slide, Inches(0.6), Inches(5.25), Inches(12), Inches(1.3), summary)


def slide_20_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, CORP_BLUE)
    _add_circle(slide, Inches(10.5), Inches(-0.5), Inches(3), CORP_BLUE_DARK)
    _add_circle(slide, Inches(-0.8), Inches(5.5), Inches(2), CORP_BLUE_DARK)

    _add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
                  '2026.04.28 重點回顧 & 後市展望', font_size=32, font_color=WHITE,
                  bold=True, alignment=PP_ALIGN.CENTER)
    _add_bg_rect(slide, Inches(5), Inches(1.3), Inches(3.3), Inches(0.04), WHITE)

    _add_rounded_rect(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.2), WHITE)
    summary_lines = [
        {'text': '本集五大重點回顧', 'size': 18, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  1. 台股衝破4萬點：4/27盤中40,194歷史新高，收39,616(+684pt)，日K收黑。外資賣500億，內資（投信+散戶）撐盤，1,300檔下跌', 'size': 13, 'color': DARK_GRAY},
        {'text': '  2. 台積電條款：買盤剛啟動。散戶先搶，投信需1-3個月（契約修改→董事會→主管機關）。預估1,500-2,000億陸續流入', 'size': 13, 'color': DARK_GRAY},
        {'text': '  3. 費半均值回歸：估值60倍，資金輪向七巨頭/軟體。標普500&那指齊創新高，YTD+5.1%，行情健康', 'size': 13, 'color': DARK_GRAY},
        {'text': '  4. 金管會五箭：亞資中心7月全台落地；虛擬資產管理條例下半年啟動；台灣全球自由度第5名', 'size': 13, 'color': DARK_GRAY},
        {'text': '  5. IPO搶錢大作戰：2026年預估1,600億歷史新高，SpaceX/Anthropic/OpenAI待下半年登場', 'size': 13, 'color': DARK_GRAY},
    ]
    _add_multiline_text_box(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(3.0),
                            summary_lines)

    _add_rounded_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.8), WHITE)
    next_lines = [
        {'text': '後市觀察重點', 'size': 18, 'color': CORP_BLUE, 'bold': True},
        {'text': '', 'size': 5, 'color': DARK_GRAY},
        {'text': '  投信台積電條款何時落地（關鍵時程）；新台幣匯率是否持續走強（外資回補前兆）', 'size': 14, 'color': CHARCOAL},
        {'text': '  七巨頭財報週結果（軟體能否接棒費半）；ETF贖回潮是否出現（系統性風險指標）', 'size': 14, 'color': CHARCOAL},
        {'text': '  IPO下半年資金胃納量；高雄亞資中心7月全台落地後實質進展', 'size': 14, 'color': CHARCOAL},
    ]
    _add_multiline_text_box(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.6),
                            next_lines)

    _add_text_box(slide, Inches(2), Inches(7.0), Inches(9.3), Inches(0.4),
                  '資料來源：早晨財經速解讀 | 分析日期：2026-04-28 | 本報告僅供參考，非投資建議',
                  font_size=11, font_color=SKY_BLUE, alignment=PP_ALIGN.CENTER)


# ============================================================
# Main Build
# ============================================================

def build_presentation():
    setup_matplotlib()
    print('Generating charts...')
    charts = gen_charts()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print('Building slides...')

    slide_01_cover(prs)
    print('  [1/20] Cover')

    slide_02_contents(prs)
    print('  [2/20] Contents')

    slide_03_part_one(prs)
    print('  [3/20] Part 1 divider')

    slide_04_us_weekly(prs, charts)
    print('  [4/20] US indices — S&P/NASDAQ new high, SOX mean reversion')

    slide_05_taiwan_market(prs, charts)
    print('  [5/20] Taiwan 40,000 milestone')

    slide_06_part_two(prs)
    print('  [6/20] Part 2 divider')

    slide_07_tsmc_clause(prs, charts)
    print('  [7/20] TSMC clause — 3-phase analysis')

    slide_08_active_etf_tsmc(prs, charts)
    print('  [8/20] Active ETF TSMC weighting')

    slide_09_fund_rotation(prs, charts)
    print('  [9/20] Fund rotation: hardware to software')

    slide_10_part_three(prs)
    print('  [10/20] Part 3 divider')

    slide_11_taiwan_economy(prs, charts)
    print('  [11/20] Taiwan economy & wealth effect')

    slide_12_fsc_policy(prs, charts)
    print('  [12/20] FSC 5-arrow policy')

    slide_13_ipo_market(prs, charts)
    print('  [13/20] IPO market 2026 record')

    slide_14_part_four(prs)
    print('  [14/20] Part 4 divider')

    slide_15_market_forces(prs, charts)
    print('  [15/20] 5 forces supporting Taiwan 40,000')

    slide_16_foreign_vs_domestic(prs, charts)
    print('  [16/20] Foreign sell vs domestic bull')

    slide_17_part_five(prs)
    print('  [17/20] Part 5 divider')

    slide_18_recommendations(prs, charts)
    print('  [18/20] Investment recommendations')

    slide_19_allocation(prs, charts)
    print('  [19/20] Asset allocation')

    slide_20_summary(prs)
    print('  [20/20] Summary & outlook')

    prs.save(OUTPUT_PPTX)
    print(f'\n✅ Saved: {OUTPUT_PPTX}')


if __name__ == '__main__':
    build_presentation()
